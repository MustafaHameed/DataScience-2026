from __future__ import annotations

from pathlib import Path
from typing import Iterable

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from docx2pdf import convert
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt

ROOT = Path(r"e:\Documents\Github\RnD\~DataScience-2026")
DOCX_PATH = ROOT / "AI_Agents_Lecture_Handout.docx"
PDF_PATH = ROOT / "AI_Agents_Lecture_Handout.pdf"
PPTX_PATH = ROOT / "AI_Agents_Lecture_Presentation.pptx"

TITLE_COLOR = RGBColor(31, 78, 121)
ACCENT_COLOR = RGBColor(47, 117, 181)
TEXT_COLOR = RGBColor(35, 35, 35)
PPT_TITLE = PptRGBColor(31, 78, 121)
PPT_ACCENT = PptRGBColor(47, 117, 181)
PPT_TEXT = PptRGBColor(35, 35, 35)

SECTIONS = [
    (
        "1. What Are AI Agents?",
        [
            "An AI agent is a system that can perceive context, reason about goals, choose actions, and use tools to complete tasks with limited human intervention.",
            "Unlike a single prompt-response chatbot, an agent can maintain state, plan multiple steps, call external systems, and adapt when conditions change.",
            "Agents are especially useful for workflows such as research assistance, report generation, customer support, coding help, and task automation.",
        ],
    ),
    (
        "2. Core Components of an AI Agent",
        [
            "Goal or instruction: defines what success looks like.",
            "Model or reasoning engine: interprets user intent and decides next steps.",
            "Memory or state: stores conversation history, constraints, or intermediate outputs.",
            "Tools and actions: APIs, databases, search, calculators, code execution, or file operations.",
            "Policy and guardrails: ensure safety, privacy, accuracy, and compliance.",
            "Feedback loop: evaluates results and decides whether to continue, revise, or stop.",
        ],
    ),
    (
        "3. Typical Agent Workflow",
        [
            "Receive a goal from the user or system.",
            "Understand the problem and break it into sub-tasks.",
            "Select the best tool or data source for each step.",
            "Execute actions, observe results, and update memory.",
            "Check whether the goal is met; if not, iterate.",
            "Return a final answer or perform the requested action.",
        ],
    ),
    (
        "4. Single-Agent vs Multi-Agent Systems",
        [
            "A single-agent system is simpler to design and easier to govern, making it suitable for focused tasks.",
            "A multi-agent system distributes responsibilities across specialized agents such as planner, researcher, critic, and executor.",
            "Multi-agent designs can improve modularity and scale, but they add coordination overhead and more failure modes.",
        ],
    ),
    (
        "5. Real-World Use Cases",
        [
            "Education: tutoring, quiz generation, assignment feedback, and personalized study planning.",
            "Business: support ticket triage, meeting summarization, market research, and workflow automation.",
            "Software engineering: bug triage, code review assistance, documentation drafting, and test generation.",
            "Healthcare and public services: guided intake, information retrieval, and administrative task support under strict human oversight.",
        ],
    ),
    (
        "6. Benefits of AI Agents",
        [
            "Reduce repetitive manual work.",
            "Operate across multiple tools and data sources.",
            "Support faster decision-making through automation.",
            "Enable scalable personalization for users and learners.",
            "Provide consistent workflows with traceable steps.",
        ],
    ),
    (
        "7. Risks and Challenges",
        [
            "Hallucinations or incorrect reasoning can lead to poor decisions.",
            "Over-automation may create hidden errors if humans are removed from review.",
            "Tool misuse can affect privacy, security, or data integrity.",
            "Prompt injection and malicious inputs can manipulate agent behavior.",
            "Poor memory management can cause outdated or inconsistent actions.",
        ],
    ),
    (
        "8. Design Best Practices",
        [
            "Keep the agent scope narrow before expanding to complex autonomy.",
            "Use structured tool interfaces and validated outputs.",
            "Log important actions for review and auditing.",
            "Add human approval checkpoints for high-impact actions.",
            "Measure quality using accuracy, latency, cost, and safety metrics.",
            "Test with normal, edge-case, and adversarial prompts.",
        ],
    ),
    (
        "9. Example Classroom Scenario",
        [
            "A student asks for help understanding a research topic.",
            "The agent clarifies the goal, searches trusted sources, drafts an explanation, and suggests readings.",
            "A verification step checks whether the explanation matches the requested level and cites the right material.",
            "The instructor reviews the output before it is shared widely.",
        ],
    ),
    (
        "10. Key Takeaways",
        [
            "AI agents combine reasoning, memory, tools, and feedback to perform multi-step tasks.",
            "They create value when tasks require coordination, repetition, and timely decisions.",
            "Strong guardrails and human oversight are essential for reliable deployment.",
            "The best agent systems are practical, measurable, and aligned with real user needs.",
        ],
    ),
]

SLIDES = [
    {
        "title": "AI Agents",
        "subtitle": "Lecture handout and presentation\nFoundations, workflows, use cases, and best practices",
    },
    {
        "title": "Learning Objectives",
        "bullets": [
            "Define what an AI agent is",
            "Explain core components and workflow",
            "Differentiate single-agent and multi-agent systems",
            "Discuss use cases, benefits, risks, and best practices",
        ],
    },
    {
        "title": "What Is an AI Agent?",
        "bullets": [
            "Perceives context and interprets goals",
            "Plans and executes multi-step actions",
            "Uses tools, memory, and feedback loops",
            "Operates with limited human intervention",
        ],
    },
    {
        "title": "Core Components",
        "bullets": [
            "Goal or task definition",
            "Reasoning model",
            "Memory and state",
            "Tools and external systems",
            "Guardrails and evaluation",
        ],
    },
    {
        "title": "Agent Workflow",
        "bullets": [
            "Receive goal",
            "Break problem into steps",
            "Select tools and act",
            "Observe results and iterate",
            "Stop when success criteria are met",
        ],
    },
    {
        "title": "Single vs Multi-Agent",
        "bullets": [
            "Single-agent: simpler and easier to govern",
            "Multi-agent: specialized roles and parallel work",
            "Trade-off: more power but more coordination complexity",
        ],
    },
    {
        "title": "Use Cases",
        "bullets": [
            "Education and tutoring",
            "Customer support and operations",
            "Research and reporting",
            "Software engineering assistance",
        ],
    },
    {
        "title": "Benefits",
        "bullets": [
            "Automation of repetitive tasks",
            "Faster workflows across tools",
            "Scalable personalization",
            "Consistent and traceable processes",
        ],
    },
    {
        "title": "Risks and Challenges",
        "bullets": [
            "Hallucinations and factual errors",
            "Prompt injection and unsafe tool use",
            "Privacy and compliance concerns",
            "Over-automation without review",
        ],
    },
    {
        "title": "Best Practices",
        "bullets": [
            "Start narrow and measure outcomes",
            "Validate tool outputs",
            "Add human approval for sensitive steps",
            "Monitor cost, latency, quality, and safety",
        ],
    },
    {
        "title": "Classroom Example",
        "bullets": [
            "Agent receives a student research query",
            "Searches trusted sources and drafts explanation",
            "Checks clarity and citation quality",
            "Instructor reviews before publication",
        ],
    },
    {
        "title": "Takeaways",
        "bullets": [
            "AI agents are practical systems for multi-step work",
            "They need memory, tools, and feedback",
            "Good governance is as important as model quality",
            "Human oversight remains essential",
        ],
    },
]


def set_document_margins(document: Document) -> None:
    for section in document.sections:
        section.top_margin = Inches(0.7)
        section.bottom_margin = Inches(0.7)
        section.left_margin = Inches(0.8)
        section.right_margin = Inches(0.8)


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shading = OxmlElement("w:shd")
    shading.set(qn("w:fill"), fill)
    tc_pr.append(shading)


def add_bullets(document: Document, items: Iterable[str]) -> None:
    for item in items:
        paragraph = document.add_paragraph(style="List Bullet")
        paragraph.paragraph_format.space_after = Pt(3)
        run = paragraph.add_run(item)
        run.font.name = "Calibri"
        run.font.size = Pt(11)
        run.font.color.rgb = TEXT_COLOR


def add_section(document: Document, title: str, points: list[str]) -> None:
    heading = document.add_paragraph()
    heading.paragraph_format.space_before = Pt(8)
    heading.paragraph_format.space_after = Pt(4)
    run = heading.add_run(title)
    run.bold = True
    run.font.name = "Calibri"
    run.font.size = Pt(14)
    run.font.color.rgb = TITLE_COLOR
    add_bullets(document, points)


def build_docx() -> Path:
    document = Document()
    set_document_margins(document)

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title.add_run("Lecture Handout: AI Agents")
    title_run.bold = True
    title_run.font.name = "Calibri"
    title_run.font.size = Pt(22)
    title_run.font.color.rgb = TITLE_COLOR

    subtitle = document.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle_run = subtitle.add_run("Concepts, Architecture, Workflows, Use Cases, and Best Practices")
    subtitle_run.italic = True
    subtitle_run.font.name = "Calibri"
    subtitle_run.font.size = Pt(11)
    subtitle_run.font.color.rgb = ACCENT_COLOR

    intro = document.add_paragraph()
    intro.paragraph_format.space_before = Pt(8)
    intro.paragraph_format.space_after = Pt(10)
    intro_run = intro.add_run(
        "This handout introduces AI agents as systems that can reason, plan, use tools, and complete multi-step tasks. "
        "It is designed for classroom teaching, quick revision, and lecture discussion."
    )
    intro_run.font.name = "Calibri"
    intro_run.font.size = Pt(11)
    intro_run.font.color.rgb = TEXT_COLOR

    table = document.add_table(rows=1, cols=2)
    table.style = "Table Grid"
    table.autofit = True
    headers = table.rows[0].cells
    headers[0].text = "Quick View"
    headers[1].text = "Details"
    set_cell_shading(headers[0], "D9EAF7")
    set_cell_shading(headers[1], "D9EAF7")
    for cell in headers:
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                run.bold = True
                run.font.name = "Calibri"
                run.font.size = Pt(11)
                run.font.color.rgb = TITLE_COLOR

    overview_rows = [
        ("Definition", "An AI agent is a goal-directed system that can observe, reason, act, and adapt."),
        ("Key Ability", "Performs multi-step tasks using models, tools, memory, and feedback."),
        ("Examples", "Research assistant, tutor, coding helper, support automation, workflow orchestrator."),
        ("Main Concern", "Reliability, safety, privacy, and the need for human oversight."),
    ]
    for left, right in overview_rows:
        row = table.add_row().cells
        row[0].text = left
        row[1].text = right

    for title_text, points in SECTIONS:
        add_section(document, title_text, points)

    closing = document.add_paragraph()
    closing.paragraph_format.space_before = Pt(10)
    closing_run = closing.add_run(
        "Discussion prompt: In which educational or organizational task would an AI agent create the most value, and what guardrails would you require before deploying it?"
    )
    closing_run.bold = True
    closing_run.font.name = "Calibri"
    closing_run.font.size = Pt(11)
    closing_run.font.color.rgb = ACCENT_COLOR

    document.save(DOCX_PATH)
    return DOCX_PATH


def build_pdf(docx_path: Path) -> Path:
    convert(str(docx_path), str(PDF_PATH))
    return PDF_PATH


def format_slide_title(shape) -> None:
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.runs[0]
    run.font.name = "Calibri"
    run.font.size = PptPt(26)
    run.font.bold = True
    run.font.color.rgb = PPT_TITLE


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PptRGBColor(248, 250, 252)

    title_shape = slide.shapes.title
    title_shape.text = title
    format_slide_title(title_shape)

    body = slide.placeholders[1]
    text_frame = body.text_frame
    text_frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.name = "Calibri"
            run.font.size = PptPt(20)
            run.font.color.rgb = PPT_TEXT

    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        PptInches(0),
        PptInches(0),
        PptInches(13.33),
        PptInches(0.3),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = PPT_ACCENT
    banner.line.fill.background()


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PptRGBColor(248, 250, 252)

    title_shape = slide.shapes.title
    title_shape.text = title
    format_slide_title(title_shape)

    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    for paragraph in subtitle_shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.name = "Calibri"
            run.font.size = PptPt(18)
            run.font.color.rgb = PPT_TEXT

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        PptInches(0.7),
        PptInches(4.8),
        PptInches(11.6),
        PptInches(0.6),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = PPT_ACCENT
    accent.line.fill.background()
    accent.text_frame.text = "AI agents = models + memory + tools + feedback + guardrails"
    para = accent.text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    for run in para.runs:
        run.font.name = "Calibri"
        run.font.size = PptPt(18)
        run.font.bold = True
        run.font.color.rgb = PptRGBColor(255, 255, 255)


def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    first = SLIDES[0]
    add_title_slide(prs, first["title"], first["subtitle"])
    for slide in SLIDES[1:]:
        add_bullet_slide(prs, slide["title"], slide["bullets"])

    prs.save(PPTX_PATH)
    return PPTX_PATH


def main() -> None:
    ROOT.mkdir(parents=True, exist_ok=True)
    docx_path = build_docx()
    pdf_path = build_pdf(docx_path)
    pptx_path = build_pptx()
    print(f"Generated DOCX: {docx_path}")
    print(f"Generated PDF: {pdf_path}")
    print(f"Generated PPTX: {pptx_path}")


if __name__ == "__main__":
    main()
