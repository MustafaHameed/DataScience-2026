"""
Generate ARM PowerPoint Presentation — Computer Accessories Dataset
===================================================================
Run this script to produce:  ARM_Presentation.pptx
Requires:  pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Colour palette ──────────────────────────────────────────────
DARK_BG   = RGBColor(0x2C, 0x3E, 0x50)
MED_BG    = RGBColor(0x34, 0x49, 0x5E)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x00, 0x00, 0x00)
BLUE      = RGBColor(0x29, 0x80, 0xB9)
GREEN     = RGBColor(0x27, 0xAE, 0x60)
ORANGE    = RGBColor(0xE6, 0x7E, 0x22)
RED       = RGBColor(0xE7, 0x4C, 0x3C)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ─── Helper functions ─────────────────────────────────────────────

def add_bg(slide, color):
    """Set solid background colour for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_slide(slide, left, top, width, height, bullets,
                     font_size=20, color=BLACK, bold_first=False,
                     font_name="Calibri", spacing=Pt(8)):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        if bold_first and ":" in bullet:
            # We can't do partial bold easily; just bold the whole line
            p.font.bold = True
    return tf

def title_slide(title, subtitle="", bg=DARK_BG, title_color=WHITE, sub_color=LIGHT_GRAY):
    """Create a title slide with coloured background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, bg)
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.5),
                title, font_size=40, bold=True, color=title_color,
                alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(3.8), Inches(11.7), Inches(1.0),
                    subtitle, font_size=24, color=sub_color,
                    alignment=PP_ALIGN.CENTER)
    return slide

def section_slide(title, bg=MED_BG):
    """Create a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, bg)
    add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.0),
                title, font_size=36, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    return slide

def content_slide(title, bg=WHITE):
    """Create a slide with a title at the top."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if bg != WHITE:
        add_bg(slide, bg)
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
                title, font_size=32, bold=True, color=DARK_BG,
                alignment=PP_ALIGN.LEFT)
    # Add a thin line under the title
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0.5), Inches(1.05),
                                    Inches(12.3), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    return slide

def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    """Add a table to the slide. data is a list of lists (rows × cols)."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE
    return table

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════
title_slide(
    "Association Rule Mining (ARM)",
    "Introduction with R — Computer Accessories Dataset\nDr. Mustafa Hameed"
)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — What is ARM?
# ═══════════════════════════════════════════════════════════════════
s = content_slide("What is Association Rule Mining?")
add_bullet_slide(s, Inches(0.6), Inches(1.3), Inches(11.5), Inches(3.0), [
    "ARM is a data mining technique that discovers interesting",
    "  relationships (associations) among items in large datasets.",
    "",
    "• Finds hidden patterns — which items are purchased together?",
    "• Classic use case: Market Basket Analysis",
    "• Applications: retail, e-commerce, healthcare, web mining",
], font_size=22, color=BLACK)
add_textbox(s, Inches(0.8), Inches(4.8), Inches(11.0), Inches(1.2),
            '"Customers who buy a laptop and a mouse also tend to buy a keyboard."\n'
            '{laptop, mouse}  →  {keyboard}',
            font_size=22, color=BLUE, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — Why Does ARM Matter?
# ═══════════════════════════════════════════════════════════════════
s = content_slide("Why Does ARM Matter?")
add_textbox(s, Inches(0.6), Inches(1.3), Inches(5.5), Inches(0.5),
            "Business Uses", font_size=24, bold=True, color=BLUE)
add_bullet_slide(s, Inches(0.6), Inches(1.9), Inches(5.5), Inches(3.5), [
    "• Product placement — nearby items",
    "• Bundle deals — laptop + mouse + keyboard",
    "• Cross-selling — recommend at checkout",
    "• Inventory — stock associated items together",
], font_size=20)
add_textbox(s, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.5),
            "Other Domains", font_size=24, bold=True, color=BLUE)
add_bullet_slide(s, Inches(7.0), Inches(1.9), Inches(5.5), Inches(3.5), [
    "• Medical co-diagnoses",
    "• Web click-stream patterns",
    "• Course enrollment patterns",
    "• Fraud detection",
], font_size=20)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — Key Definitions
# ═══════════════════════════════════════════════════════════════════
s = section_slide("Key Definitions")

s2 = content_slide("Key Definitions")
data = [
    ["Term", "Definition"],
    ["Transaction", "A single purchase event (one receipt / one basket)"],
    ["Item", "An individual product in a transaction"],
    ["Itemset", "A collection of one or more items, e.g. {laptop, mouse}"],
    ["Association Rule", "An implication: {A} → {B} ('if A then B')"],
    ["Antecedent (LHS)", "The 'if' part — {laptop, mouse}"],
    ["Consequent (RHS)", "The 'then' part — {keyboard}"],
]
add_table(s2, Inches(0.6), Inches(1.3), Inches(12.0), Inches(4.5),
          len(data), 2, data,
          col_widths=[Inches(3.5), Inches(8.5)])

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — The Four Key Measures
# ═══════════════════════════════════════════════════════════════════
s = content_slide("The Four Key Measures of ARM")
data = [
    ["#", "Measure", "Formula", "Question it Answers"],
    ["1", "Support", "freq(A ∪ B) / N", "How common is this itemset?"],
    ["2", "Confidence", "Supp(A ∪ B) / Supp(A)", "How reliable is this rule?"],
    ["3", "Lift", "Conf(A→B) / Supp(B)", "Is it better than random chance?"],
    ["4", "Conviction", "(1−Supp(B)) / (1−Conf)", "How dependent is the rule?"],
]
add_table(s, Inches(0.6), Inches(1.3), Inches(12.0), Inches(3.5),
          len(data), 4, data,
          col_widths=[Inches(0.6), Inches(2.2), Inches(4.2), Inches(5.0)])
add_textbox(s, Inches(0.8), Inches(5.2), Inches(11.0), Inches(1.0),
            "Support = frequency  •  Confidence = reliability  •  Lift = strength  •  Conviction = dependency",
            font_size=20, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — Dataset
# ═══════════════════════════════════════════════════════════════════
s = content_slide("Our Dataset — Computer Accessories (N = 10)")
data = [
    ["TID", "Items Purchased"],
    ["T1", "laptop, mouse, keyboard"],
    ["T2", "laptop, mouse"],
    ["T3", "laptop, keyboard, USB hub"],
    ["T4", "mouse, keyboard, mousepad"],
    ["T5", "laptop, mouse, keyboard, USB hub"],
    ["T6", "mouse, mousepad"],
    ["T7", "laptop, mouse, keyboard, webcam"],
    ["T8", "keyboard, headset"],
    ["T9", "laptop, mouse, keyboard, mousepad"],
    ["T10", "laptop, keyboard"],
]
add_table(s, Inches(1.5), Inches(1.3), Inches(10.0), Inches(5.5),
          len(data), 2, data,
          col_widths=[Inches(1.5), Inches(8.5)])

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — Item Counts
# ═══════════════════════════════════════════════════════════════════
s = content_slide("Item Counts (by hand)")
data = [
    ["Item", "Transactions", "Count", "Support"],
    ["laptop", "T1,T2,T3,T5,T7,T9,T10", "7", "0.70"],
    ["mouse", "T1,T2,T4,T5,T6,T7,T9", "7", "0.70"],
    ["keyboard", "T1,T3,T4,T5,T7,T8,T9,T10", "8", "0.80"],
    ["USB hub", "T3, T5", "2", "0.20"],
    ["mousepad", "T4, T6, T9", "3", "0.30"],
    ["webcam", "T7", "1", "0.10"],
    ["headset", "T8", "1", "0.10"],
]
add_table(s, Inches(0.6), Inches(1.3), Inches(12.0), Inches(4.5),
          len(data), 4, data,
          col_widths=[Inches(2.0), Inches(5.0), Inches(1.5), Inches(3.5)])

# ═══════════════════════════════════════════════════════════════════
# SLIDES 8-11 — Measure definitions + step-by-step calculations
# ═══════════════════════════════════════════════════════════════════

# --- SUPPORT ---
s = section_slide("Measure 1 — Support")
s = content_slide("Support — Definition")
add_textbox(s, Inches(0.6), Inches(1.3), Inches(12.0), Inches(1.5),
            "Support measures how frequently an itemset appears in the dataset.\n\n"
            "Support(X) = (# transactions containing X) / N",
            font_size=24, color=BLACK)
add_bullet_slide(s, Inches(0.6), Inches(3.2), Inches(12.0), Inches(3.0), [
    "• Higher support = more common pattern",
    "• Used to filter out rare patterns (set a minimum threshold)",
    "• Support of 0.50 = itemset appears in 50% of all transactions",
], font_size=22)

s = content_slide("Support — Step-by-Step: {laptop, mouse} → {keyboard}")
add_bullet_slide(s, Inches(0.6), Inches(1.3), Inches(12.0), Inches(4.5), [
    "Step 1: Find transactions containing {laptop, mouse, keyboard}:",
    "    T1: laptop, mouse, keyboard  ✓",
    "    T5: laptop, mouse, keyboard, USB hub  ✓",
    "    T7: laptop, mouse, keyboard, webcam  ✓",
    "    T9: laptop, mouse, keyboard, mousepad  ✓",
    "",
    "Step 2: Count = 4  out of  N = 10",
    "",
    "Support({laptop, mouse, keyboard}) = 4 / 10 = 0.40",
    "",
    "✅ Interpretation: 40% of all transactions contain laptop + mouse + keyboard.",
], font_size=20)

# --- CONFIDENCE ---
s = section_slide("Measure 2 — Confidence")
s = content_slide("Confidence — Definition")
add_textbox(s, Inches(0.6), Inches(1.3), Inches(12.0), Inches(1.5),
            "Confidence measures how often the rule is correct.\n"
            "Given that A is purchased, how often is B also purchased?\n\n"
            "Confidence(A → B) = Support(A ∪ B) / Support(A)",
            font_size=24, color=BLACK)
add_bullet_slide(s, Inches(0.6), Inches(3.5), Inches(12.0), Inches(2.5), [
    "• Ranges from 0 to 1 (0% to 100%)",
    "• Higher is better, but always verify with Lift!",
    "• Confidence of 0.80 = '80% of the time when A is bought, B is also bought'",
], font_size=22)

s = content_slide("Confidence — Step-by-Step: {laptop, mouse} → {keyboard}")
add_bullet_slide(s, Inches(0.6), Inches(1.3), Inches(12.0), Inches(4.5), [
    "Step 1: Support({laptop, mouse, keyboard}) = 4/10 = 0.40",
    "",
    "Step 2: Transactions with {laptop, mouse}: T1, T2, T5, T7, T9 → 5",
    "    Support({laptop, mouse}) = 5/10 = 0.50",
    "",
    "Step 3: Confidence = 0.40 / 0.50 = 0.80",
    "",
    "✅ Interpretation: When a customer buys laptop + mouse,",
    "   there is an 80% chance they also buy a keyboard.",
], font_size=22)

# --- LIFT ---
s = section_slide("Measure 3 — Lift")
s = content_slide("Lift — Definition")
add_textbox(s, Inches(0.6), Inches(1.3), Inches(12.0), Inches(1.5),
            "Lift measures whether A and B appear together more often\n"
            "than expected by chance.\n\n"
            "Lift(A → B) = Confidence(A → B) / Support(B)",
            font_size=24, color=BLACK)
data = [
    ["Lift Value", "Meaning"],
    ["Lift > 1", "Positive association — items appear together more than expected"],
    ["Lift = 1", "No association — items are independent"],
    ["Lift < 1", "Negative association — items avoid each other"],
]
add_table(s, Inches(1.0), Inches(3.8), Inches(11.0), Inches(2.5),
          len(data), 2, data,
          col_widths=[Inches(2.5), Inches(8.5)])

s = content_slide("Lift — Step-by-Step: {laptop, mouse} → {keyboard}")
add_bullet_slide(s, Inches(0.6), Inches(1.3), Inches(12.0), Inches(4.5), [
    "Step 1: Confidence = 0.80  (from previous)",
    "",
    "Step 2: Support({keyboard}) = 8/10 = 0.80",
    "",
    "Step 3: Lift = 0.80 / 0.80 = 1.00",
    "",
    "⚠️ Interpretation: Lift = 1.00 means laptop + mouse does NOT",
    "   increase the chance of buying keyboard — they are INDEPENDENT.",
    "",
    "   This happens because keyboard is already in 80% of transactions!",
], font_size=22)

# --- CONVICTION ---
s = section_slide("Measure 4 — Conviction")
s = content_slide("Conviction — Definition")
add_textbox(s, Inches(0.6), Inches(1.3), Inches(12.0), Inches(1.5),
            "Conviction measures the degree of dependence.\n"
            "It compares the expected error rate (if independent) to actual.\n\n"
            "Conviction(A → B) = (1 − Support(B)) / (1 − Confidence(A → B))",
            font_size=24, color=BLACK)
data = [
    ["Conviction", "Meaning"],
    ["Conv = 1", "Items are independent (no dependency)"],
    ["Conv > 1", "Rule is meaningful — items are dependent"],
    ["Conv → ∞", "Perfect rule (confidence = 100%)"],
]
add_table(s, Inches(1.0), Inches(3.8), Inches(11.0), Inches(2.0),
          len(data), 2, data,
          col_widths=[Inches(2.5), Inches(8.5)])

s = content_slide("Conviction — Step-by-Step: {laptop, mouse} → {keyboard}")
add_bullet_slide(s, Inches(0.6), Inches(1.3), Inches(12.0), Inches(4.0), [
    "Step 1: 1 − Support({keyboard}) = 1 − 0.80 = 0.20",
    "",
    "Step 2: 1 − Confidence = 1 − 0.80 = 0.20",
    "",
    "Step 3: Conviction = 0.20 / 0.20 = 1.00",
    "",
    "⚠️ Interpretation: Conviction = 1.00 confirms INDEPENDENCE",
    "   — consistent with Lift = 1.00.",
], font_size=22)

# ═══════════════════════════════════════════════════════════════════
# SLIDE — Rule 2: {mouse} → {mousepad}
# ═══════════════════════════════════════════════════════════════════
s = content_slide("Another Rule: {mouse} → {mousepad}")
data = [
    ["Step", "Calculation", "Result"],
    ["Txns with {mouse, mousepad}", "T4, T6, T9", "3"],
    ["Support({mouse, mousepad})", "3/10", "0.30"],
    ["Support({mouse})", "7/10", "0.70"],
    ["Support({mousepad})", "3/10", "0.30"],
    ["Confidence", "0.30 / 0.70", "0.429"],
    ["Lift", "0.429 / 0.30", "1.43"],
    ["Conviction", "(1−0.30)/(1−0.429)", "1.23"],
]
add_table(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(4.5),
          len(data), 3, data,
          col_widths=[Inches(4.0), Inches(4.0), Inches(3.5)])
add_textbox(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(1.0),
            "✅ Lift = 1.43, Conviction = 1.23 — a meaningful positive association!\n"
            "Buying a mouse increases the chance of buying a mousepad by 43% over random.",
            font_size=20, bold=True, color=GREEN)

# ═══════════════════════════════════════════════════════════════════
# SLIDE — Rule 3: {laptop} → {mouse}
# ═══════════════════════════════════════════════════════════════════
s = content_slide("Another Rule: {laptop} → {mouse}")
data = [
    ["Step", "Calculation", "Result"],
    ["Txns with {laptop, mouse}", "T1,T2,T5,T7,T9", "5"],
    ["Support({laptop, mouse})", "5/10", "0.50"],
    ["Support({laptop})", "7/10", "0.70"],
    ["Support({mouse})", "7/10", "0.70"],
    ["Confidence", "0.50 / 0.70", "0.714"],
    ["Lift", "0.714 / 0.70", "1.02"],
    ["Conviction", "(1−0.70)/(1−0.714)", "1.05"],
]
add_table(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(4.5),
          len(data), 3, data,
          col_widths=[Inches(4.0), Inches(4.0), Inches(3.5)])
add_textbox(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(1.0),
            "Lift = 1.02, Conviction = 1.05 — very slight positive association. Nearly independent.",
            font_size=20, color=ORANGE)

# ═══════════════════════════════════════════════════════════════════
# SLIDE — Comparing All Three Rules
# ═══════════════════════════════════════════════════════════════════
s = content_slide("Comparing All Three Rules")
data = [
    ["Rule", "Support", "Confidence", "Lift", "Conviction", "Verdict"],
    ["{laptop,mouse}→{keyboard}", "0.40", "0.80", "1.00", "1.00", "Independent"],
    ["{mouse}→{mousepad}", "0.30", "0.43", "1.43", "1.23", "✅ Useful"],
    ["{laptop}→{mouse}", "0.50", "0.71", "1.02", "1.05", "Borderline"],
]
add_table(s, Inches(0.4), Inches(1.3), Inches(12.5), Inches(2.5),
          len(data), 6, data,
          col_widths=[Inches(3.5), Inches(1.5), Inches(2.0), Inches(1.5), Inches(2.0), Inches(2.0)])
add_textbox(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.0),
            "Key Insight: The rule with the highest confidence (0.80) is the least useful (Lift=1.0).\n"
            "The rule with the highest lift (1.43) is the most actionable!",
            font_size=22, bold=True, color=RED)

# ═══════════════════════════════════════════════════════════════════
# SLIDE — How to Interpret Results
# ═══════════════════════════════════════════════════════════════════
s = content_slide("How to Interpret Results")
data = [
    ["Measure", "What to Look For"],
    ["Support", "Filter out rare patterns (min 0.10–0.30)"],
    ["Confidence", "High confidence (≥ 0.50) = rule fires correctly often"],
    ["Lift", "> 1 = useful  •  > 1.5 = strong  •  > 2 = very strong"],
    ["Conviction", "> 1 = dependent  •  > 1.5 = strong dependency"],
]
add_table(s, Inches(0.6), Inches(1.3), Inches(12.0), Inches(3.0),
          len(data), 2, data,
          col_widths=[Inches(2.5), Inches(9.5)])
add_textbox(s, Inches(0.8), Inches(4.8), Inches(11.5), Inches(2.0),
            "⚠️ WATCH OUT: High confidence + low lift → the consequent is just very popular!\n"
            "Always check Lift and Conviction before acting on a rule!",
            font_size=22, bold=True, color=RED)

# ═══════════════════════════════════════════════════════════════════
# SLIDE — Business Actions
# ═══════════════════════════════════════════════════════════════════
s = content_slide("Business Actions Based on Lift")
data = [
    ["Lift Range", "Action"],
    ["Lift > 2.0", "Strong bundle — create a package deal"],
    ["Lift 1.5–2.0", "Cross-sell — recommend at checkout"],
    ["Lift 1.1–1.5", "Place nearby — adjacent shelves"],
    ["Lift ≈ 1.0", "No action — items are independent"],
    ["Lift < 1.0", "Separate — avoid placing together"],
]
add_table(s, Inches(1.5), Inches(1.3), Inches(10.0), Inches(3.5),
          len(data), 2, data,
          col_widths=[Inches(3.0), Inches(7.0)])

# ═══════════════════════════════════════════════════════════════════
# SLIDE — The Apriori Algorithm
# ═══════════════════════════════════════════════════════════════════
s = content_slide("The Apriori Algorithm")
add_bullet_slide(s, Inches(0.6), Inches(1.3), Inches(12.0), Inches(4.5), [
    "1. Set minimum support — prune rare itemsets early",
    "2. Find frequent 1-itemsets — {laptop}, {mouse}, {keyboard}, …",
    "3. Build frequent 2-itemsets by joining 1-itemsets",
    "4. Continue growing until no new frequent itemsets found",
    "5. Generate rules that meet minimum confidence",
    "",
    "Anti-monotone property:",
    "A subset of a frequent itemset must also be frequent.",
    "If {laptop, mouse} is not frequent, then {laptop, mouse, keyboard}",
    "cannot be frequent — so we PRUNE it!",
], font_size=22)

# ═══════════════════════════════════════════════════════════════════
# SLIDE — Summary of Four Measures
# ═══════════════════════════════════════════════════════════════════
s = section_slide("Summary of the Four Measures")

s = content_slide("Four Measures — Final Summary")
data = [
    ["Measure", "Formula", "Good Values"],
    ["Support", "freq(A∪B) / N", "≥ 0.10"],
    ["Confidence", "Supp(A∪B) / Supp(A)", "≥ 0.50"],
    ["Lift", "Conf / Supp(B)", "> 1 (higher = better)"],
    ["Conviction", "(1−Supp(B))/(1−Conf)", "> 1 (higher = stronger)"],
]
add_table(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(3.0),
          len(data), 3, data,
          col_widths=[Inches(2.5), Inches(5.0), Inches(4.0)])

# ═══════════════════════════════════════════════════════════════════
# SLIDE — Key Takeaways
# ═══════════════════════════════════════════════════════════════════
s = section_slide("Key Takeaways")

s2 = content_slide("Key Takeaways")
add_bullet_slide(s2, Inches(0.6), Inches(1.3), Inches(12.0), Inches(5.0), [
    "• Support — how common is the pattern?",
    "• Confidence — how reliable is the rule?",
    "• Lift > 1 — the rule is better than random guessing",
    "• Conviction > 1 — items are genuinely dependent",
    "• Always check Lift before trusting a high-confidence rule",
    "• The Apriori algorithm prunes efficiently (anti-monotone property)",
    "• arules + arulesViz make ARM easy in R",
    "",
    "Try It Yourself!",
    "  1. Change supp and conf thresholds — how does the rule count change?",
    "  2. Add more transactions — do the measures improve?",
    "  3. Filter for 'mouse' on the RHS — what predicts mouse purchases?",
], font_size=22)

# ═══════════════════════════════════════════════════════════════════
# SLIDE — References
# ═══════════════════════════════════════════════════════════════════
s = content_slide("References")
add_bullet_slide(s, Inches(0.6), Inches(1.3), Inches(12.0), Inches(5.0), [
    "1. Agrawal, R., & Srikant, R. (1994). Fast algorithms for mining",
    "   association rules. VLDB.",
    "2. Hahsler, M., Grün, B., & Hornik, K. (2005). arules — A computational",
    "   environment for mining association rules. JSS, 14(15).",
    "3. Hahsler, M. (2017). arulesViz: Visualizing Association Rules.",
    "   JSS, 76(2).",
    "4. Brin, S., Motwani, R., Ullman, J., & Tsur, S. (1997). Dynamic itemset",
    "   counting and implication rules. ACM SIGMOD.",
], font_size=18)

# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "ARM_Presentation.pptx")
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
