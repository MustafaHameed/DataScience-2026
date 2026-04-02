"""
Generate mtcars EDA PowerPoint Presentation
============================================
Run:  python generate_mtcars_pptx.py
Produces:  mtcars_EDA_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Colour palette ──────────────────────────────────────────────
DARK_BG    = RGBColor(0x2C, 0x3E, 0x50)
MED_BG     = RGBColor(0x34, 0x49, 0x5E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
BLUE       = RGBColor(0x29, 0x80, 0xB9)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)
RED        = RGBColor(0xE7, 0x4C, 0x3C)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ─── Helpers ──────────────────────────────────────────────────────
def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def tb(slide, l, t, w, h, text, sz=18, bold=False, color=BLACK,
       align=PP_ALIGN.LEFT, name="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = name; p.alignment = align
    return tf

def bullets(slide, l, t, w, h, items, sz=20, color=BLACK, spacing=Pt(6)):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.name = "Calibri"; p.space_after = spacing
    return tf

def title_slide(title, subtitle=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK_BG)
    tb(s, Inches(.8), Inches(2), Inches(11.7), Inches(1.5),
       title, 40, True, WHITE, PP_ALIGN.CENTER)
    if subtitle:
        tb(s, Inches(.8), Inches(3.8), Inches(11.7), Inches(1),
           subtitle, 24, False, LIGHT_GRAY, PP_ALIGN.CENTER)
    return s

def section(title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, MED_BG)
    tb(s, Inches(.8), Inches(2.5), Inches(11.7), Inches(2),
       title, 36, True, WHITE, PP_ALIGN.CENTER)
    return s

def content(title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb(s, Inches(.5), Inches(.3), Inches(12.3), Inches(.8),
       title, 32, True, DARK_BG)
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(.5), Inches(1.05), Inches(12.3), Pt(3))
    sh.fill.solid(); sh.fill.fore_color.rgb = BLUE
    sh.line.fill.background()
    return s

def add_tbl(slide, l, t, w, h, data, cw=None):
    rows, cols = len(data), len(data[0])
    tsh = slide.shapes.add_table(rows, cols, l, t, w, h)
    tbl = tsh.table
    if cw:
        for i, ww in enumerate(cw):
            tbl.columns[i].width = ww
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c); cell.text = str(data[r][c])
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14); p.font.name = "Calibri"
                if r == 0:
                    p.font.bold = True; p.font.color.rgb = WHITE
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE
    return tbl

# ═════════════════════════════════════════════════════════════════
# SLIDES
# ═════════════════════════════════════════════════════════════════

# 1 — Title
title_slide("Exploratory Data Analysis (EDA)",
            "The mtcars Dataset — Introduction with R\nDr. Mustafa Hameed")

# 2 — What is EDA?
s = content("What is Exploratory Data Analysis?")
bullets(s, Inches(.6), Inches(1.3), Inches(11.5), Inches(3), [
    "EDA is the FIRST step in any data science project.",
    "",
    "• Inspect — what does the data look like?",
    "• Summarize — what are the key statistics?",
    "• Visualize — what patterns exist?",
    "• Clean — any missing values or errors?",
    "",
    "Goal: understand your data BEFORE building models.",
], sz=22)

# 3 — Why EDA Matters
s = content("Why EDA Matters")
add_tbl(s, Inches(.6), Inches(1.3), Inches(12), Inches(3.5),
    [["Purpose", "What You Learn"],
     ["Catch errors early", "Missing values, wrong types, impossible values"],
     ["Understand distributions", "Skewed? Outliers? Normal?"],
     ["Find relationships", "Which variables are correlated?"],
     ["Guide modeling", "Choose the right algorithm"],
     ["Communicate", "Visualizations tell the story"]],
    cw=[Inches(3.5), Inches(8.5)])

# 4 — Key Statistical Concepts
s = content("Key Statistical Concepts")
add_tbl(s, Inches(.6), Inches(1.3), Inches(12), Inches(4),
    [["Term", "Definition"],
     ["Mean", "Average: sum of values ÷ count"],
     ["Median", "Middle value when data is sorted"],
     ["Std Deviation", "How spread out values are around the mean"],
     ["Quartiles", "Q1 (25%), Q2 (50% = median), Q3 (75%)"],
     ["Correlation", "Linear relationship strength: −1 to +1"],
     ["Outlier", "A value unusually far from the rest"]],
    cw=[Inches(2.5), Inches(9.5)])

# 5 — The mtcars Dataset
s = content("The mtcars Dataset")
bullets(s, Inches(.6), Inches(1.3), Inches(12), Inches(2), [
    "Source: 1974 Motor Trend US magazine",
    "32 automobiles (1973-74 models), 11 variables",
], sz=22)
add_tbl(s, Inches(.6), Inches(2.8), Inches(12), Inches(4.2),
    [["Variable", "Description", "Type"],
     ["mpg", "Miles per gallon", "Continuous"],
     ["cyl", "Cylinders (4, 6, 8)", "Categorical"],
     ["disp", "Displacement (cu.in.)", "Continuous"],
     ["hp", "Horsepower", "Continuous"],
     ["drat", "Rear axle ratio", "Continuous"],
     ["wt", "Weight (×1000 lbs)", "Continuous"],
     ["qsec", "Quarter-mile time (sec)", "Continuous"],
     ["vs", "Engine (0=V, 1=straight)", "Binary"],
     ["am", "Trans (0=auto, 1=manual)", "Binary"],
     ["gear", "Forward gears (3,4,5)", "Categorical"],
     ["carb", "Carburetors", "Categorical"]],
    cw=[Inches(1.8), Inches(6.2), Inches(4)])

# 6 — Section: Step 1
section("Step 1 — Load & Inspect")

s = content("Load and Inspect the Data")
tb(s, Inches(.6), Inches(1.3), Inches(12), Inches(.6),
   "R Code:", 20, True, BLUE)
bullets(s, Inches(.8), Inches(1.9), Inches(11), Inches(2.5), [
    "dim(mtcars)          # 32 rows × 11 columns",
    "head(mtcars)         # First 6 rows",
    "str(mtcars)          # Variable names, types, samples",
    "names(mtcars)        # Column names only",
], sz=18, color=DARK_BG)
tb(s, Inches(.6), Inches(4.2), Inches(12), Inches(.6),
   "What to look for:", 20, True, BLUE)
bullets(s, Inches(.8), Inches(4.8), Inches(11), Inches(2), [
    "• 32 observations (cars), 11 features",
    "• All variables stored as 'num' — some are actually categorical",
    "• No missing values in this dataset",
], sz=20)

# 7 — Section: Step 2
section("Step 2 — Summary Statistics")

s = content("Summary Statistics")
tb(s, Inches(.6), Inches(1.3), Inches(12), Inches(.4),
   "summary(mtcars) gives Min, Q1, Median, Mean, Q3, Max for each variable.", 20)
add_tbl(s, Inches(.6), Inches(2), Inches(12), Inches(2.5),
    [["Variable", "Min", "Mean", "Median", "Max", "SD"],
     ["mpg", "10.40", "20.09", "19.20", "33.90", "6.03"],
     ["hp", "52.0", "146.7", "123.0", "335.0", "68.6"],
     ["wt", "1.51", "3.22", "3.33", "5.42", "0.98"],
     ["disp", "71.1", "230.7", "196.3", "472.0", "123.9"]],
    cw=[Inches(2), Inches(1.5), Inches(2), Inches(2), Inches(2), Inches(2.5)])
tb(s, Inches(.6), Inches(5), Inches(12), Inches(1.5),
   "Key insight: MPG ranges from 10.4 to 33.9 — wide variation.\n"
   "HP has SD = 68.6 — horsepower varies enormously across these cars.",
   20, color=GREEN)

# 8 — Step 3: Frequency
s = content("Frequency Counts — Categorical Variables")
add_tbl(s, Inches(.6), Inches(1.3), Inches(5.5), Inches(2.5),
    [["Cylinders", "Count"],
     ["4", "11"],
     ["6", "7"],
     ["8", "14"]],
    cw=[Inches(2.5), Inches(3)])
add_tbl(s, Inches(6.8), Inches(1.3), Inches(5.5), Inches(2.5),
    [["Transmission", "Count"],
     ["Automatic (0)", "19"],
     ["Manual (1)", "13"]],
    cw=[Inches(3), Inches(2.5)])
add_tbl(s, Inches(.6), Inches(4.2), Inches(5.5), Inches(2.5),
    [["Gears", "Count"],
     ["3", "15"],
     ["4", "12"],
     ["5", "5"]],
    cw=[Inches(2.5), Inches(3)])
tb(s, Inches(6.8), Inches(4.5), Inches(5.5), Inches(2),
   "Most cars: 8 cylinders, automatic, 3 gears.\n"
   "This tells us the 'typical' car in this 1974 sample.",
   20, color=BLUE)

# 9 — Section: Visualizations
section("Step 3 — Visualizations")

# 10 — Histograms
s = content("Histograms — Understand Distributions")
bullets(s, Inches(.6), Inches(1.3), Inches(12), Inches(1.5), [
    "A histogram groups data into bins and shows frequency.",
    "Use: hist(mtcars$mpg, col='steelblue', breaks=10)",
], sz=20)
tb(s, Inches(.6), Inches(2.8), Inches(5.5), Inches(.5),
   "MPG Distribution", 22, True, BLUE)
bullets(s, Inches(.6), Inches(3.4), Inches(5.5), Inches(2.5), [
    "• Most cars: 15–22 MPG",
    "• Slightly right-skewed",
    "• Mean ≈ 20.1 MPG",
    "• A few outliers above 30 MPG",
], sz=18)
tb(s, Inches(7), Inches(2.8), Inches(5.5), Inches(.5),
   "HP Distribution", 22, True, BLUE)
bullets(s, Inches(7), Inches(3.4), Inches(5.5), Inches(2.5), [
    "• Most cars: 50–175 HP",
    "• Right-skewed",
    "• Mean ≈ 146.7 HP",
    "• A few high-performance outliers (> 250 HP)",
], sz=18)

# 11 — Boxplots
s = content("Boxplots — Compare Groups")
bullets(s, Inches(.6), Inches(1.3), Inches(12), Inches(1.2), [
    "A boxplot shows the median, quartiles, and outliers.",
    "Use: boxplot(mpg ~ cyl, data = mtcars)",
], sz=20)
tb(s, Inches(.6), Inches(2.6), Inches(5.5), Inches(.5),
   "MPG by Cylinders", 22, True, BLUE)
bullets(s, Inches(.6), Inches(3.2), Inches(5.5), Inches(3), [
    "• 4-cyl: highest MPG (~26.7 avg)",
    "• 6-cyl: moderate (~19.7 avg)",
    "• 8-cyl: lowest MPG (~15.1 avg)",
    "• Clear negative relationship:",
    "  more cylinders → worse fuel economy",
], sz=18)
tb(s, Inches(7), Inches(2.6), Inches(5.5), Inches(.5),
   "MPG by Transmission", 22, True, BLUE)
bullets(s, Inches(7), Inches(3.2), Inches(5.5), Inches(3), [
    "• Automatic: avg 17.1 MPG",
    "• Manual: avg 24.4 MPG",
    "• Manual cars tend to be",
    "  more fuel-efficient",
    "• (Confounded with weight/size)",
], sz=18)

# 12 — Scatter Plots
s = content("Scatter Plots — Relationships Between Variables")
bullets(s, Inches(.6), Inches(1.3), Inches(12), Inches(1), [
    "Each dot = one car. Add a regression line with abline(lm(...)).",
], sz=20)
tb(s, Inches(.6), Inches(2.3), Inches(5.5), Inches(.5),
   "MPG vs Weight", 22, True, BLUE)
bullets(s, Inches(.6), Inches(2.9), Inches(5.5), Inches(3), [
    "• Strong negative relationship",
    "• r = −0.87",
    "• Heavier cars → worse fuel economy",
    "• Nearly linear trend",
    "• Weight is the #1 predictor of MPG",
], sz=18)
tb(s, Inches(7), Inches(2.3), Inches(5.5), Inches(.5),
   "MPG vs Horsepower", 22, True, BLUE)
bullets(s, Inches(7), Inches(2.9), Inches(5.5), Inches(3), [
    "• Strong negative relationship",
    "• r = −0.78",
    "• More HP → worse fuel economy",
    "• Slightly curved (non-linear)",
], sz=18)

# 13 — Pairs Plot
s = content("Pairs Plot — Multiple Relationships at Once")
bullets(s, Inches(.6), Inches(1.3), Inches(12), Inches(5), [
    "pairs(mtcars[, c('mpg','hp','wt','disp','qsec')])",
    "",
    "A pairs plot shows scatter plots for EVERY pair of variables.",
    "",
    "What to look for:",
    "• Strong diagonal trends → strong correlation",
    "• Cloud-like scatter → weak/no relationship",
    "• Curved patterns → non-linear relationship",
    "",
    "Key findings from the pairs plot:",
    "• mpg vs wt: strong negative (linear)",
    "• mpg vs disp: strong negative",
    "• wt vs disp: strong positive (big engines → heavy cars)",
    "• qsec vs hp: moderate negative (more HP → faster quarter mile)",
], sz=20)

# 14 — Section: Correlation
section("Step 4 — Correlation Analysis")

s = content("Correlation Matrix")
add_tbl(s, Inches(.6), Inches(1.3), Inches(12), Inches(4),
    [["", "mpg", "hp", "wt", "disp", "drat", "qsec"],
     ["mpg", "1.00", "−0.78", "−0.87", "−0.85", "0.68", "0.42"],
     ["hp", "−0.78", "1.00", "0.66", "0.79", "−0.45", "−0.71"],
     ["wt", "−0.87", "0.66", "1.00", "0.89", "−0.71", "−0.17"],
     ["disp", "−0.85", "0.79", "0.89", "1.00", "−0.71", "−0.43"],
     ["drat", "0.68", "−0.45", "−0.71", "−0.71", "1.00", "0.09"],
     ["qsec", "0.42", "−0.71", "−0.17", "−0.43", "0.09", "1.00"]],
    cw=[Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5),
        Inches(2), Inches(2), Inches(2)])
add_tbl(s, Inches(1), Inches(5.5), Inches(11), Inches(1.5),
    [["Range", "Strength"],
     ["±0.70 to ±1.00", "Strong"],
     ["±0.40 to ±0.69", "Moderate"],
     ["±0.10 to ±0.39", "Weak"]],
    cw=[Inches(4), Inches(7)])

# 15 — Strongest Correlations
s = content("Key Correlations — What They Mean")
add_tbl(s, Inches(.6), Inches(1.3), Inches(12), Inches(4),
    [["Variable Pair", "r", "Strength", "Interpretation"],
     ["mpg vs wt", "−0.87", "Strong −", "Heavier cars use more fuel"],
     ["mpg vs disp", "−0.85", "Strong −", "Bigger engines use more fuel"],
     ["mpg vs hp", "−0.78", "Strong −", "More powerful cars use more fuel"],
     ["wt vs disp", "+0.89", "Strong +", "Heavier cars have bigger engines"],
     ["mpg vs drat", "+0.68", "Moderate +", "Higher axle ratio → better MPG"],
     ["hp vs qsec", "−0.71", "Strong −", "More HP → faster quarter mile"]],
    cw=[Inches(2.5), Inches(1.2), Inches(2), Inches(6.3)])

# 16 — Section: Subsetting
section("Step 5 — Subsetting & Aggregation")

s = content("Filtering and Aggregating Data")
tb(s, Inches(.6), Inches(1.3), Inches(12), Inches(.5),
   "Filter: Fuel-efficient cars (MPG > 20)", 22, True, BLUE)
bullets(s, Inches(.8), Inches(1.9), Inches(11), Inches(1.5), [
    "efficient <- mtcars[mtcars$mpg > 20, c('mpg','cyl','hp','wt')]",
    "Result: 14 out of 32 cars qualify",
], sz=18, color=DARK_BG)
tb(s, Inches(.6), Inches(3.5), Inches(12), Inches(.5),
   "Aggregate: Mean MPG by Cylinder Group", 22, True, BLUE)
add_tbl(s, Inches(.8), Inches(4.2), Inches(5), Inches(2),
    [["Cylinders", "Mean MPG"],
     ["4", "26.66"],
     ["6", "19.74"],
     ["8", "15.10"]],
    cw=[Inches(2.5), Inches(2.5)])
add_tbl(s, Inches(6.5), Inches(4.2), Inches(5.5), Inches(1.5),
    [["Transmission", "Mean MPG"],
     ["Automatic", "17.15"],
     ["Manual", "24.39"]],
    cw=[Inches(3), Inches(2.5)])

# 17 — Summary Findings
s = content("Summary of Key Findings")
add_tbl(s, Inches(.6), Inches(1.3), Inches(12), Inches(3.5),
    [["Finding", "Evidence"],
     ["MPG varies widely", "Range: 10.4 – 33.9"],
     ["Heavier cars → worse MPG", "r = −0.87 (strong negative)"],
     ["More cylinders → worse MPG", "4-cyl: 26.7, 8-cyl: 15.1"],
     ["Manual > Automatic for MPG", "Manual: 24.4, Auto: 17.1"],
     ["HP and displacement hurt MPG", "r = −0.78 and −0.85"]],
    cw=[Inches(4.5), Inches(7.5)])

# 18 — EDA Cheat Sheet
s = content("EDA Cheat Sheet — R Functions")
add_tbl(s, Inches(.4), Inches(1.3), Inches(12.5), Inches(5.5),
    [["Task", "Function", "Example"],
     ["View data", "head(), tail()", "head(mtcars)"],
     ["Dimensions", "dim(), nrow(), ncol()", "dim(mtcars)"],
     ["Structure", "str()", "str(mtcars)"],
     ["Summary", "summary(), mean(), sd()", "summary(mtcars$mpg)"],
     ["Frequency", "table()", "table(mtcars$cyl)"],
     ["Histogram", "hist()", "hist(mtcars$mpg)"],
     ["Boxplot", "boxplot()", "boxplot(mpg ~ cyl, data=mtcars)"],
     ["Scatter", "plot()", "plot(mtcars$wt, mtcars$mpg)"],
     ["Correlation", "cor()", "cor(mtcars$mpg, mtcars$wt)"],
     ["Filter", "[ ] or subset()", "mtcars[mtcars$mpg > 20, ]"],
     ["Aggregate", "aggregate()", "aggregate(mpg~cyl, mtcars, mean)"]],
    cw=[Inches(2.2), Inches(4), Inches(6.3)])

# 19 — Key Takeaways
section("Key Takeaways")

s = content("Key Takeaways")
bullets(s, Inches(.6), Inches(1.3), Inches(12), Inches(5), [
    "1. EDA is the FIRST step — always explore before modeling",
    "2. Use head(), str(), summary() to inspect data",
    "3. Histograms show distributions — look for skew and outliers",
    "4. Boxplots compare groups — great for categorical variables",
    "5. Scatter plots reveal relationships — add regression lines",
    "6. Correlation quantifies linear relationships (−1 to +1)",
    "7. Subsetting focuses analysis on interesting subgroups",
    "8. Aggregation summarizes by group — compare means across categories",
    "",
    "For mtcars:",
    "• Weight is the strongest predictor of fuel economy",
    "• More cylinders, more HP, larger engines → worse MPG",
    "• Manual cars tend to be more efficient (but confounded!)",
], sz=22)

# 20 — Practice
s = content("Practice Exercises")
bullets(s, Inches(.6), Inches(1.3), Inches(12), Inches(5), [
    "1. Create a histogram of wt (weight). Is it skewed?",
    "2. Make a boxplot of hp grouped by cyl. What do you observe?",
    "3. Calculate cor(mtcars$hp, mtcars$wt). Interpret the result.",
    "4. Filter cars with HP > 200. How many? What is their avg MPG?",
    "5. Create scatter plot of disp vs mpg with regression line.",
    "6. Which 4-cylinder car has the best MPG?",
    "7. Is the relationship between qsec and hp linear?",
], sz=22)

# 21 — References
s = content("References")
bullets(s, Inches(.6), Inches(1.3), Inches(12), Inches(4), [
    "1. Henderson, H. V., & Velleman, P. F. (1981). Building multiple",
    "   regression models interactively. Biometrics, 37, 391–411.",
    "2. R Core Team (2024). R: A language and environment for",
    "   statistical computing.",
    "3. Tukey, J. W. (1977). Exploratory Data Analysis. Addison-Wesley.",
    "4. Wickham, H. (2016). ggplot2: Elegant Graphics for Data",
    "   Analysis. Springer.",
], sz=18)

# ─── Save ──────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "mtcars_EDA_Presentation.pptx")
prs.save(out)
print(f"✅ Saved: {out}")
print(f"   Slides: {len(prs.slides)}")
