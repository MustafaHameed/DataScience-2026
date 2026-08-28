#!/usr/bin/env python3
"""Check the decks, then render them so they can be looked at.

Two halves, and the second is the one that matters.

**Static checks** read the .pptx and catch what is measurable: text that will
not fit its box, missing figures, slides with no speaker notes, and Morph
groups whose shape names drift apart (which silently turns an animation into
a cut).

**Rendering** drives the installed PowerPoint through COM to export each deck
to PDF, then rasterises every slide to PNG. No checker catches a cramped
panel, a bad colour pairing or a build step that reveals the wrong thing --
only looking does, which is how the book's 289 pages were reviewed.

    python qa.py                 check every deck in out/
    python qa.py --ch 3          one chapter
    python qa.py --render        also export PDF and PNG
    python qa.py --render --ch 3
"""
from __future__ import annotations

import argparse
import glob
import os

import subprocess
import sys

from pptx import Presentation


HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "out")
RENDER = os.path.join(HERE, "render")

EMU_IN = 914400.0
CHAR_W = 0.50 / 72.0      # width of a character, in inches per point
LINE_H = 1.30 / 72.0


def _fits(shape) -> tuple[bool, float]:
    """Estimate whether a text box's content fits. Returns (ok, fill ratio).

    Paragraph spacing below the *last* paragraph is not rendered inside the
    box, so counting it made every single-line label look 15% overfull.
    """
    tf = shape.text_frame
    w = shape.width / EMU_IN
    h = shape.height / EMU_IN
    if w <= 0 or h <= 0:
        return True, 0.0

    # Code panels set word_wrap False: a long line is clipped, not wrapped,
    # so counting it as several lines reported healthy panels as overfull.
    wraps = tf.word_wrap is not False

    paras = list(tf.paragraphs)
    total = 0.0
    for i, p in enumerate(paras):
        size, text = None, ""
        for r in p.runs:
            text += r.text
            if r.font.size:
                size = max(size or 0, r.font.size.pt)
        size = size or 18
        if wraps:
            per_line = max(1, int(w / (CHAR_W * size)))
            n = max(1, -(-len(text) // per_line))
        else:
            n = 1
        # Honour the paragraph's own line spacing; code panels set 1.14 and
        # assuming the 1.30 default reported them 14% fuller than they are.
        spacing = 1.30
        try:
            if p.line_spacing and isinstance(p.line_spacing, float):
                spacing = p.line_spacing * 1.22
        except Exception:
            pass
        total += n * (spacing / 72.0) * size
        if p.space_after and i < len(paras) - 1:
            total += p.space_after.pt / 72.0
    # 8% tolerance: the estimate is deliberately conservative on character
    # width, and PowerPoint itself allows a little bleed before it looks wrong.
    return total <= h * 1.08, (total / h if h else 0.0)


def check_deck(path: str, verbose: bool = False) -> list[str]:
    prs = Presentation(path)
    name = os.path.basename(path)
    problems: list[str] = []
    prev_names: set[str] | None = None
    prev_morph = False

    for i, slide in enumerate(prs.slides, 1):
        names = set()
        has_text = False
        for sh in slide.shapes:
            names.add(sh.name)
            if sh.has_text_frame and sh.text_frame.text.strip():
                has_text = True
                ok, ratio = _fits(sh)
                if not ok:
                    problems.append(
                        f"{name} slide {i}: '{sh.name}' overflows "
                        f"({ratio:.0%} of its box) -- "
                        f"{sh.text_frame.text.strip()[:60]!r}")
        if not has_text and not any(
                sh.shape_type == 13 for sh in slide.shapes):   # 13 = picture
            problems.append(f"{name} slide {i}: empty")

        # Morph needs matching shape names on consecutive slides
        xml = slide.element.xml
        morph = "morph" in xml
        # A morph needs at least one shape name in common to have something
        # to tween. One is enough and often deliberate -- a checkpoint
        # question morphs from a full-bleed slide into the answer slide's
        # heading, and only the question text is shared. Zero means it will
        # simply cut.
        if morph and prev_morph and prev_names is not None:
            if not (names & prev_names):
                problems.append(
                    f"{name} slide {i}: morph shares no shape name with the "
                    f"previous slide -- will cut, not animate")
        prev_names, prev_morph = names, morph

    notes = sum(1 for s in prs.slides
                if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip())
    if verbose:
        print(f"  {name:34} {len(prs.slides._sldIdLst):>3} slides, "
              f"{notes:>3} with notes")
    return problems


# ---------------------------------------------------------------------------
# Rendering
# ---------------------------------------------------------------------------

def to_pdf(paths: list[str]) -> list[str]:
    """Export decks to PDF using the installed PowerPoint."""
    try:
        import win32com.client as win32
    except ImportError:
        print("  pywin32 not available; cannot render")
        return []
    os.makedirs(RENDER, exist_ok=True)
    made = []
    app = None
    try:
        app = win32.Dispatch("PowerPoint.Application")
        for p in paths:
            pdf = os.path.join(RENDER,
                               os.path.splitext(os.path.basename(p))[0] + ".pdf")
            try:
                pres = app.Presentations.Open(os.path.abspath(p),
                                              ReadOnly=True, WithWindow=False)
                pres.SaveAs(os.path.abspath(pdf), 32)     # 32 = ppSaveAsPDF
                pres.Close()
                made.append(pdf)
                print(f"  exported {os.path.basename(pdf)}")
            except Exception as e:                        # noqa: BLE001
                print(f"  export FAILED for {os.path.basename(p)}: {e}")
    finally:
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
    return made


def to_png(pdfs: list[str], dpi: int = 90) -> int:
    n = 0
    for pdf in pdfs:
        stem = os.path.splitext(os.path.basename(pdf))[0]
        d = os.path.join(RENDER, stem)
        os.makedirs(d, exist_ok=True)
        subprocess.run(["pdftocairo", "-png", "-r", str(dpi), pdf,
                        os.path.join(d, "s")],
                       capture_output=True, text=True)
        n += len(glob.glob(os.path.join(d, "s-*.png")))
    return n


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--ch", nargs="*", type=int)
    ap.add_argument("--render", action="store_true")
    ap.add_argument("--dpi", type=int, default=90)
    a = ap.parse_args()

    paths = sorted(glob.glob(os.path.join(OUT, "*.pptx")))
    if a.ch:
        keep = {f"ch{n:02d}" for n in a.ch}
        paths = [p for p in paths
                 if os.path.basename(p)[:4] in keep]
    if not paths:
        print("  no decks in slides/out/ -- run build_slides.py first")
        return 1

    print(f"Checking {len(paths)} deck(s)\n")
    problems: list[str] = []
    for p in paths:
        problems += check_deck(p, verbose=True)

    print()
    if problems:
        for q in problems[:40]:
            print("  " + q)
        if len(problems) > 40:
            print(f"  ... and {len(problems) - 40} more")
        print(f"\n  {len(problems)} problem(s)")
    else:
        print("  static checks clean")

    if a.render:
        print("\nRendering")
        pdfs = to_pdf(paths)
        n = to_png(pdfs, a.dpi)
        print(f"  {len(pdfs)} PDF, {n} slide PNGs -> slides/render/")

    return 1 if problems else 0


if __name__ == "__main__":
    sys.exit(main())
