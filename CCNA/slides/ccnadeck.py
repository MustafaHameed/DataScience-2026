#!/usr/bin/env python3
"""The design system for the CCNA lecture decks.

One place that knows what a slide looks like, so 34 decks look like one course
rather than 34 attempts. ``build_slides.py`` describes *what* goes on a slide;
this module decides *how* it is drawn.

Two things here are worth knowing before editing:

**Morph.** python-pptx cannot create entrance animations. Progressive builds are
therefore made of consecutive slides that carry the same shapes, in the same
places, under the same *names* -- PowerPoint's Morph transition then tweens
between them by itself. ``Deck.step()`` opens such a group and every shape it
draws gets a stable name. Rename a shape between steps and the morph breaks
into a cut, which is why ``qa.py`` checks the names match.

**Fitting.** There is no reliable auto-shrink in python-pptx: PowerPoint only
recalculates autofit when it opens the file. Sizes are therefore chosen up
front by ``fit_pt()`` from the amount of text, and ``qa.py`` renders every
slide to PNG afterwards so overflow is caught by looking rather than by hoping.
"""
from __future__ import annotations

import os

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
# One asset directory for the volume and the decks, so the title page of
# the book and the title slide of a deck cannot end up with different
# versions of the same mark.
ASSETS = os.path.join(os.path.dirname(HERE), "assets")

# ---------------------------------------------------------------------------
# Palette -- lifted verbatim from ccna.sty so slides and book match exactly
# ---------------------------------------------------------------------------

PRIMARY = RGBColor(0x1A, 0x52, 0x76)
SECONDARY = RGBColor(0x2E, 0x86, 0xC1)
ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
NEUTRAL = RGBColor(0x70, 0x7B, 0x7C)
LIGHTBG = RGBColor(0xEB, 0xF5, 0xFB)
DEFBG = RGBColor(0xFD, 0xF2, 0xE9)
GREENHL = RGBColor(0x27, 0xAE, 0x60)
GOLDHL = RGBColor(0xB7, 0x95, 0x0B)
CODEBG = RGBColor(0x1B, 0x26, 0x31)      # darker than the book's, for projection
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEARBLACK = RGBColor(0x1C, 0x25, 0x2B)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x8A, 0x95, 0x9A)

DOMAIN = {
    "1": RGBColor(0x1F, 0x61, 0x8D),
    "2": RGBColor(0x11, 0x78, 0x64),
    "3": RGBColor(0x6C, 0x34, 0x83),
    "4": RGBColor(0x92, 0x2B, 0x21),
    "5": RGBColor(0xB7, 0x95, 0x0B),
}
DOMAIN_NAME = {
    "1": "Network Infrastructure and Connectivity",
    "2": "Switching and Network Access",
    "3": "IP Routing",
    "4": "Network Services and Security",
    "5": "AI, Network Operations and Management",
}

BODY_FONT = "Segoe UI"
MONO_FONT = "Consolas"

# ---------------------------------------------------------------------------
# Geometry (inches). One grid, used by every layout.
# ---------------------------------------------------------------------------

W, H = 13.333, 7.5
ML, MR = 0.75, 0.75
MT, MB = 0.52, 0.42
CW = W - ML - MR                 # 11.83 content width
TITLE_Y, TITLE_H = 0.50, 0.82
RULE_Y = 1.40
BODY_Y = 1.66
BODY_H = 5.32                    # down to 6.98
FOOT_Y = 7.02


def _rgb(c: RGBColor) -> RGBColor:
    return c


# ---------------------------------------------------------------------------
# Text fitting
# ---------------------------------------------------------------------------

def fit_pt(text: str, width_in: float, height_in: float,
           ideal: int = 20, floor: int = 13) -> int:
    """Pick a point size that should fit *text* into the given box.

    Rough but stable: at N pt a Segoe UI character is about 0.50*N/72 in wide
    and a line is about 1.30*N/72 in tall. Start at *ideal* and step down until
    the estimate fits, never below *floor* -- below that a lecture slide is
    unreadable from the back of a room and the content should be split instead.
    """
    text = text or ""
    for pt in range(ideal, floor - 1, -1):
        cw = 0.50 * pt / 72.0
        lh = 1.30 * pt / 72.0
        per_line = max(1, int(width_in / cw))
        lines = 0
        for para in text.split("\n"):
            lines += max(1, -(-len(para) // per_line))
        if lines * lh <= height_in:
            return pt
    return floor


# Per-character widths in em, for the few places where a box is sized to one
# short string rather than to wrapped prose. The flat 0.50 em that fit_pt uses
# is a fair average for a sentence but far too narrow for an all-capital term:
# it made the "DHCPDISCOVER" key-term chip wrap onto a second line.
_NARROW = set("iljtfrI.,;:'!|()[]{} •")
_WIDE = set("mwMW@%")


def text_w(text: str, size: int) -> float:
    """Width of a short single-line string in inches, erring generous."""
    em = 0.0
    for c in text:
        if c in _NARROW:
            em += 0.30
        elif c in _WIDE:
            em += 0.88
        elif c.isupper() or c.isdigit():
            em += 0.62
        else:
            em += 0.51
    return em * size / 72.0


def _hang(p, marker: str, size: int) -> None:
    """Give a paragraph a hanging indent the width of its marker.

    The bullet is an inline run here, not a PowerPoint list marker, so nothing
    indents the wrapped lines by itself: a two-line item put its second line
    hard against the left margin and read as a separate paragraph. ``marL``
    sets where wrapped lines start; the matching negative ``indent`` pulls the
    first line back out to the margin so the marker still hangs.
    """
    mar = int(text_w(marker, size) * 914400)     # EMU
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(mar))
    pPr.set("indent", str(-mar))


def _ground(c: RGBColor) -> RGBColor:
    """The hue darkened enough to carry white text on a full-bleed slide.

    Four of the five domain colours sit near luminance 80. The fifth, the gold
    of domain 5, is at 143, so a title slide in it left white type washed out
    and swallowed the gold ribbon on the department shield, while the other
    four were fine. Scaling to a common luminance fixes that one deck without
    inventing a colour or special-casing it by name.

    Nothing is ever lightened: a colour already dark enough comes back exactly
    as the book defines it, so slides and volume still match.
    """
    lum = (c[0] * 299 + c[1] * 587 + c[2] * 114) / 1000
    if lum <= 88:
        return c
    k = 88 / lum
    return RGBColor(*(min(255, int(v * k)) for v in c))


class Deck:
    """A single chapter's deck."""

    def __init__(self, chapter: int, title: str, topics: list[str],
                 subtitle: str = "", author: str = "Dr. Mustafa Hameed",
                 dept: str = "Department of Information Technology",
                 course: str = "CCNA 200-301 v2.0"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(W)
        self.prs.slide_height = Inches(H)
        self.chapter = chapter
        self.title = title
        self.topics = topics
        self.subtitle = subtitle
        self.author = author
        self.dept = dept
        self.course = course
        self.hue = DOMAIN.get(topics[0].split(".")[0], PRIMARY) if topics \
            else PRIMARY
        self.ground = _ground(self.hue)
        self._blank = self.prs.slide_layouts[6]
        self._step_group: str | None = None
        self._step_n = 0

    # -- slide plumbing ----------------------------------------------------

    def _new(self, bg: RGBColor | None = None):
        s = self.prs.slides.add_slide(self._blank)
        if bg is not None:
            f = s.background.fill
            f.solid()
            f.fore_color.rgb = bg
        return s

    def _transition(self, slide, kind: str = "fade", ms: int = 500) -> None:
        """Attach a slide transition by injecting the XML python-pptx omits."""
        sld = slide.element
        for old in sld.findall(qn("p:transition")):
            sld.remove(old)
        ns = ('xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility'
              '/2006" xmlns:p="http://schemas.openxmlformats.org/presentationml'
              '/2006/main" xmlns:p14="http://schemas.microsoft.com/office/'
              'powerpoint/2010/main"')
        if kind == "morph":
            xml = (
                f'<mc:AlternateContent {ns}>'
                f'<mc:Choice Requires="p14">'
                f'<p:transition spd="med" p14:dur="{ms}">'
                f'<p159:morph xmlns:p159="http://schemas.microsoft.com/office/'
                f'powerpoint/2015/09/main" option="byObject"/>'
                f'</p:transition></mc:Choice>'
                f'<mc:Fallback><p:transition spd="med"><p:fade/></p:transition>'
                f'</mc:Fallback></mc:AlternateContent>')
        else:
            xml = (f'<mc:AlternateContent {ns}><mc:Choice Requires="p14">'
                   f'<p:transition spd="med" p14:dur="{ms}"><p:fade/>'
                   f'</p:transition></mc:Choice><mc:Fallback>'
                   f'<p:transition spd="med"><p:fade/></p:transition>'
                   f'</mc:Fallback></mc:AlternateContent>')
        try:
            from lxml import etree
            frag = etree.fromstring(xml)
            cs = sld.find(qn("p:cSld"))
            cs.addnext(frag) if cs is not None else sld.append(frag)
        except Exception:
            pass          # a missing transition is cosmetic, never fatal

    def notes(self, slide, text: str) -> None:
        if not text:
            return
        slide.notes_slide.notes_text_frame.text = text[:4000]

    # -- primitives --------------------------------------------------------

    def rect(self, slide, x, y, w, h, fill=None, line=None, name=None,
             lw=1.0, shape=MSO_SHAPE.RECTANGLE):
        sh = slide.shapes.add_shape(shape, Inches(x), Inches(y),
                                    Inches(w), Inches(h))
        if fill is None:
            sh.fill.background()
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = fill
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = line
            sh.line.width = Pt(lw)
        sh.shadow.inherit = False
        if name:
            sh.name = name
        sh.text_frame.text = ""
        return sh

    def text(self, slide, x, y, w, h, runs, size=20, color=NEARBLACK,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=None, name=None, space_after=6, line_spacing=1.0,
             wrap=True):
        """Place styled runs. *runs* is a str, a list of Run, or a list of those."""
        box = slide.shapes.add_textbox(Inches(x), Inches(y),
                                       Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = wrap
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)

        paras = runs if isinstance(runs, list) and runs and \
            isinstance(runs[0], list) else [runs]
        for i, para in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_after = Pt(space_after)
            p.line_spacing = line_spacing
            self._runs_into(p, para, size, color, bold, font)
        if name:
            box.name = name
        return box

    def _runs_into(self, p, para, size, color, bold, font):
        if isinstance(para, str):
            para = [_Plain(para)]
        for r in para:
            chip = getattr(r, "chip", None)
            if chip:
                run = p.add_run()
                run.text = f" {chip} "
                run.font.size = Pt(max(10, size - 4))
                run.font.bold = True
                run.font.name = BODY_FONT
                run.font.color.rgb = DOMAIN.get(chip.split(".")[0], PRIMARY)
                continue
            run = p.add_run()
            run.text = r.text
            f = run.font
            f.size = Pt(size)
            f.bold = bool(bold or getattr(r, "bold", False))
            f.italic = bool(getattr(r, "italic", False))
            mono = getattr(r, "mono", False)
            f.name = font or (MONO_FONT if mono else BODY_FONT)
            if getattr(r, "term", False):
                f.color.rgb = self.hue
                f.bold = True
            elif mono:
                f.color.rgb = PRIMARY if color == NEARBLACK else color
            else:
                f.color.rgb = color
            if mono:
                f.size = Pt(max(10, size - 1))

    # -- chrome ------------------------------------------------------------

    def _chrome(self, slide, title: str, kicker: str = "") -> None:
        """Title, domain rule and footer. Named for Morph stability."""
        if kicker:
            self.text(slide, ML, TITLE_Y - 0.04, CW, 0.26, kicker,
                      size=12, color=MUTED, bold=True, name="kicker")
            ty = TITLE_Y + 0.20
        else:
            ty = TITLE_Y
        self.text(slide, ML, ty, CW, TITLE_H, title,
                  size=fit_pt(title, CW, 0.86, ideal=32, floor=20),
                  color=PRIMARY, bold=True, name="title")
        self.rect(slide, ML, RULE_Y, CW, 0.035, fill=self.hue, name="rule")
        self._footer(slide)

    def _footer(self, slide) -> None:
        self.text(slide, ML, FOOT_Y, CW * 0.72, 0.28,
                  f"{self.course}   ·   Chapter {self.chapter}. {self.title}",
                  size=10, color=MUTED, name="footL")
        self.text(slide, ML + CW * 0.72, FOOT_Y, CW * 0.28, 0.28,
                  "  ".join(self.topics) if self.topics else "",
                  size=10, color=MUTED, align=PP_ALIGN.RIGHT, name="footR")

    def _lockup(self, slide, x, y, w, dark=False) -> None:
        """The institution in words. The marks themselves are _logos()."""
        col = WHITE if dark else PRIMARY
        sub = RGBColor(0xC8, 0xD6, 0xE5) if dark else MUTED
        self.text(slide, x, y, w, 0.30,
                  "THE ISLAMIA UNIVERSITY OF BAHAWALPUR",
                  size=12, color=col, bold=True, align=PP_ALIGN.RIGHT,
                  name="lockup1")
        self.text(slide, x, y + 0.26, w, 0.28,
                  "Department of Information Technology",
                  size=11, color=sub, align=PP_ALIGN.RIGHT, name="lockup2")

    def _logos(self, slide, right: float, y: float, h: float = 0.86,
               dark: bool = False, marks: int = 2) -> float:
        """The two institutional marks, right-aligned, ending at *right*.

        On a dark ground the university crest uses its reversed version: the
        supplied artwork is a single navy, which disappears against every one
        of the five domain colours. The department shield is multi-colour with
        a white field of its own and reads correctly on either ground.

        Returns the left edge of the row, or *right* if no artwork is present
        -- the decks were built and reviewed for months without these files
        and must keep working if they are removed.
        """
        names = ("iub-logo-white.png" if dark else "iub-logo.png",
                 "dit-logo.png")[:marks]
        gap = 0.30
        found = []
        for n in names:
            p = os.path.join(ASSETS, n)
            if os.path.exists(p):
                found.append(p)
        if not found:
            return right
        # Measure first so the row can be right-aligned in one pass.
        widths = []
        for p in found:
            with Image.open(p) as im:
                widths.append(h * im.size[0] / im.size[1])
        total = sum(widths) + gap * (len(widths) - 1)
        lx = right - total
        for p, wd in zip(found, widths):
            pic = slide.shapes.add_picture(p, Inches(lx), Inches(y),
                                           height=Inches(h))
            pic.name = "logo-" + os.path.basename(p)[:-4]
            lx += wd + gap
        return right - total

    def _chips(self, slide, x, y, topics, dark=False) -> None:
        """Blueprint topic tags, coloured by domain, as in the book."""
        cx = x
        for t in topics:
            hue = DOMAIN.get(t.split(".")[0], PRIMARY)
            w = 0.62 + 0.10 * max(0, len(t) - 3)
            self.rect(slide, cx, y, w, 0.30,
                      fill=WHITE if dark else hue,
                      line=None, name=f"chip-{t}")
            self.text(slide, cx, y + 0.035, w, 0.24, t, size=13,
                      color=hue if dark else WHITE, bold=True,
                      align=PP_ALIGN.CENTER, name=f"chiptx-{t}")
            cx += w + 0.14

    # ---------------------------------------------------------------- slides

    def title_slide(self, part_name: str = "") -> None:
        s = self._new(self.ground)
        self.rect(s, 0, 0, W, H, fill=self.ground, name="bg")
        self.rect(s, 0, H - 0.75, W, 0.75,
                  fill=RGBColor(*[max(0, c - 22) for c in self.ground]),
                  name="bgfoot")

        self.text(s, ML, 1.28, CW, 0.34, self.course.upper(),
                  size=15, color=RGBColor(0xC8, 0xD6, 0xE5), bold=True,
                  name="course")
        self.text(s, ML, 1.72, 2.2, 1.5, f"{self.chapter:02d}",
                  size=76, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True,
                  name="chnum")
        self.text(s, ML + 1.9, 1.86, CW - 1.9, 1.5, self.title,
                  size=fit_pt(self.title, CW - 2.0, 1.4, ideal=40, floor=26),
                  color=WHITE, bold=True, name="chtitle")
        if part_name:
            self.text(s, ML + 1.92, 3.30, CW - 1.9, 0.32, part_name,
                      size=15, color=RGBColor(0xC8, 0xD6, 0xE5), name="part")
        if self.topics:
            self.text(s, ML + 1.92, 3.72, 2.0, 0.26, "EXAM TOPICS",
                      size=11, color=RGBColor(0xB8, 0xC9, 0xDA), bold=True,
                      name="tlabel")
            self._chips(s, ML + 1.92, 4.00, self.topics, dark=True)

        # Left is the person, right is the institution. The department used
        # to appear on both sides.
        self.text(s, ML, 5.62, CW * 0.6, 0.30, "LECTURER", size=11,
                  color=RGBColor(0xB8, 0xC9, 0xDA), bold=True, name="authlab")
        self.text(s, ML, 5.90, CW * 0.6, 0.36, self.author,
                  size=18, color=WHITE, bold=True, name="author")
        # One institutional block, not two weak ones: the marks sit directly
        # above the words they belong to, the whole thing right-aligned to the
        # margin. At the top of the slide they had to be small enough to clear
        # a long chapter title, which left the DIT ribbon unreadable.
        self._logos(s, W - MR, 4.86, h=1.16, dark=True)
        self._lockup(s, W - MR - 4.8, 6.16, 4.8, dark=True)
        self._transition(s, "fade")
        self.notes(s, f"Chapter {self.chapter}: {self.title}. "
                      f"Blueprint topics: {', '.join(self.topics) or 'none'}.")

    def roadmap_slide(self, parts, current: int, note: str = "") -> None:
        """The six parts of the course, with the one we are in lit.

        A student meeting chapter 24 in week ten has no way to tell where it
        sits unless they are shown. Six bars, one lit, at the front of every
        deck: cheap to render and the only navigation a lecture has.
        """
        s = self._new(WHITE)
        self._chrome(s, "Where this chapter sits", "THE COURSE")
        y = BODY_Y + 0.30
        for i, (lo, hi, name) in enumerate(parts, 1):
            live = i == current
            h = 0.62
            self.rect(s, ML, y, CW, h,
                      fill=LIGHTBG if live else RGBColor(0xF7, 0xF9, 0xFA),
                      line=None, name=f"part{i}bg")
            self.rect(s, ML, y, 0.07, h,
                      fill=self.hue if live else RGBColor(0xDD, 0xE3, 0xE7),
                      name=f"part{i}rule")
            label = name.split("—")[0].strip()
            rest = name.split("—", 1)[1].strip() if "—" in name else ""
            self.text(s, ML + 0.28, y + 0.16, 1.5, 0.30, label,
                      size=14, bold=True,
                      color=self.hue if live else MUTED, name=f"part{i}n")
            self.text(s, ML + 1.65, y + 0.16, CW - 3.4, 0.30, rest,
                      size=14, bold=live,
                      color=NEARBLACK if live else MUTED, name=f"part{i}t")
            self.text(s, ML + CW - 1.7, y + 0.16, 1.5, 0.30,
                      f"chapters {lo}–{hi}", size=12,
                      color=self.hue if live else MUTED,
                      align=PP_ALIGN.RIGHT, name=f"part{i}c")
            y += h + 0.10
        self._footer(s)
        self._transition(s, "fade")
        self.notes(s, note)

    def section_slide(self, heading: str, note: str = "") -> None:
        s = self._new(WHITE)
        self.rect(s, 0, 2.55, W, 2.0, fill=self.ground, name="band")
        # Leave the heading room for the crest at the right-hand end of the
        # band rather than letting a long one run underneath it.
        self.text(s, ML, 2.92, CW - 1.5, 1.3, heading,
                  size=fit_pt(heading, CW - 1.6, 1.2, ideal=34, floor=22),
                  color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE,
                  name="secheading")
        self._logos(s, W - MR, 2.92, h=1.26, dark=True, marks=1)
        self._footer(s)
        self._transition(s, "fade")
        self.notes(s, note)

    def content_slide(self, title, paras=None, bullets=None, kicker="",
                      note="", numbered=False, tint=None, icon_text="",
                      tail=None, start=1):
        s = self._new(WHITE)
        self._chrome(s, title, kicker)
        y = BODY_Y
        inset = 0.32 if tint is not None else 0.0
        tint_shape = None
        if tint is not None:
            # Drawn now so it sits behind the text, but resized at the end
            # once the content height is known -- a three-line definition in
            # a full-height panel of colour looks like a rendering fault.
            tint_shape = self.rect(s, ML, y, CW, BODY_H, fill=tint,
                                   name="tint")
            y += 0.24
        if icon_text:
            self.text(s, ML + inset, y, CW - 2 * inset, 0.32, icon_text,
                      size=13, color=self.hue, bold=True, name="icon")
            y += 0.40

        body_text = ""
        if paras:
            body_text += "\n".join(_flat(p) for p in paras)
        if bullets:
            body_text += "\n" + "\n".join(_flat(b) for b in bullets)
        if tail:
            body_text += "\n" + "\n".join(_flat(t) for t in tail)
        # fit_pt alone under-counts a bulleted list: it misses the marker and
        # the space between items. Step down until the real estimate fits.
        room = BODY_H - (y - BODY_Y) - 0.2
        size = fit_pt(body_text, CW - 2 * inset - 0.3, room,
                      ideal=20, floor=13)
        while size > 13:
            est = _est_h(list(paras or []), CW - 2 * inset, size) + \
                  _est_h(list(bullets or []), CW - 2 * inset - 0.32, size) + \
                  _est_h(list(tail or []), CW - 2 * inset, size) + \
                  0.13 * len(bullets or [])
            if est <= room:
                break
            size -= 1

        # A small nudge for short content, and no more. Optically centring it
        # made the "(continued)" half of a split list start two inches lower
        # than the first half, so clicking between them jumped the text.
        if tint is None:
            est = _est_h(list(paras or []) + list(bullets or []),
                         CW - 2 * inset, size)
            room = BODY_H - (y - BODY_Y)
            if est < room * 0.60:
                y += min(0.42, (room - est) / 2.8)

        if paras:
            self.text(s, ML + inset, y, CW - 2 * inset, BODY_H - (y - BODY_Y),
                      list(paras), size=size, name="body", space_after=10)
            y += 0.1 + _est_h(paras, CW - 2 * inset, size)
        if bullets:
            self._bullets(s, ML + inset, y, CW - 2 * inset,
                          BODY_H - (y - BODY_Y), bullets, size, numbered,
                          start=start)
            y += _est_h(bullets, CW - 2 * inset - 0.32, size) + \
                0.13 * len(bullets)
        if tail:
            # The book's closing sentence belongs under the list it closes,
            # not above it as an introduction. Size it against the space that
            # is actually left, which after a long list may be less than the
            # body size the rest of the slide is using.
            left = max(0.42, BODY_Y + BODY_H - (y + 0.10) - 0.06)
            tsize = min(size, fit_pt("\n".join(_flat(t) for t in tail),
                                     CW - 2 * inset, left,
                                     ideal=size, floor=11))
            self.text(s, ML + inset, y + 0.10, CW - 2 * inset, left,
                      list(tail), size=tsize, name="tailbody", space_after=6)
            y += _est_h(list(tail), CW - 2 * inset, tsize) + 0.10

        if tint_shape is not None:
            used = max(1.05, min(BODY_H, y - BODY_Y + 0.30))
            tint_shape.height = Inches(used)
        self._transition(s, "fade")
        self.notes(s, note)
        return s

    def _bullets(self, slide, x, y, w, h, bullets, size, numbered=False,
                 name="bullets", upto=None, start=1):
        """Bulleted list. *upto* dims items beyond it, for stepped builds."""
        box = slide.shapes.add_textbox(Inches(x), Inches(y),
                                       Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(9)
            p.line_spacing = 1.02
            marker = f"{i + start}.  " if numbered else "•  "
            _hang(p, marker, size)
            run = p.add_run()
            run.text = marker
            run.font.size = Pt(size)
            run.font.bold = True
            run.font.name = BODY_FONT
            dim = upto is not None and i >= upto
            run.font.color.rgb = RGBColor(0xDD, 0xE2, 0xE6) if dim else self.hue
            self._runs_into(p, b, size,
                            RGBColor(0xDD, 0xE2, 0xE6) if dim else NEARBLACK,
                            False, None)
        box.name = name
        return box

    def terms_slide(self, terms: list[str], note: str = ""):
        """Key terms as a flowing field of chips.

        A vocabulary list is not a sequence of arguments, so bullets are the
        wrong form for it -- and fifteen bulleted terms overflow the body
        anyway. Laid out as chips it reads as a glossary and fits.
        """
        s = self._new(WHITE)
        self._chrome(s, "Vocabulary for this chapter", "KEY TERMS")
        pad_x, pad_y, gap = 0.18, 0.11, 0.14
        size = 16 if len(terms) <= 12 else (15 if len(terms) <= 18 else 13)
        ch = 0.30 + 2 * pad_y
        x, y = ML, BODY_Y + 0.10
        rows = 1
        for t in terms:
            w = text_w(t, size) + 2 * pad_x
            if x + w > ML + CW:
                x = ML
                y += ch + gap
                rows += 1
            self.rect(s, x, y, w, ch, fill=LIGHTBG, line=None,
                      name=f"term-{t[:18]}")
            self.rect(s, x, y, 0.035, ch, fill=self.hue, name=f"tb-{t[:18]}")
            self.text(s, x + pad_x, y + pad_y - 0.01, w - 2 * pad_x,
                      0.30, t, size=size, color=NEARBLACK,
                      name=f"tt-{t[:18]}", wrap=False)
            x += w + gap
        self.text(s, ML, min(6.86, y + ch + 0.30), CW, 0.30,
                  "Each is defined where it first matters — not here.",
                  size=13, color=MUTED, name="termnote")
        self._transition(s, "fade")
        self.notes(s, note)
        return s

    def terminal_slide(self, title, code, caption="", kicker="", note="",
                       device="", highlight=None):
        """Dark code panel. *highlight* dims every line except those listed."""
        s = self._new(WHITE)
        self._chrome(s, title, kicker)
        y = BODY_Y
        if caption:
            self.text(s, ML, y, CW, 0.34, caption, size=15, color=NEUTRAL,
                      name="caption")
            y += 0.44
        avail = BODY_H - (y - BODY_Y) - 0.05
        bar = 0.38
        lines = code.split("\n")
        size = fit_pt("\n".join(lines), CW - 0.6, avail - bar - 0.2,
                      ideal=16, floor=9)
        # Fit the panel to the code. A six-line show output in a five-inch
        # black box reads as a rendering fault, not as restraint.
        needed = bar + 0.30 + len(lines) * 1.24 * size / 72.0
        h = max(1.15, min(avail, needed))
        # Nudge a short panel down, but only a little: content that
        # starts at a different height on every slide makes the eye
        # re-find it each time. Hanging from the title rule is calmer
        # than optical centring when 40 slides go past in an hour.
        if h < avail - 0.6:
            y += min(0.42, (avail - h) / 2.6)

        self.rect(s, ML, y, CW, h, fill=CODEBG, name="codebg")
        self.rect(s, ML, y, CW, bar, fill=RGBColor(0x27, 0x35, 0x43),
                  name="codebar")
        self.text(s, ML + 0.22, y + 0.06, CW - 0.44, 0.26,
                  device or "IOS", size=12,
                  color=RGBColor(0x9F, 0xB3, 0xC8), bold=True, name="codettl")
        box = s.shapes.add_textbox(Inches(ML + 0.22), Inches(y + bar + 0.10),
                                   Inches(CW - 0.44), Inches(h - bar - 0.16))
        tf = box.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(0)
            p.line_spacing = 1.0
            run = p.add_run()
            run.text = ln or " "
            run.font.size = Pt(size)
            run.font.name = MONO_FONT
            lit = highlight is None or i in highlight
            if ln.strip().startswith("!"):
                run.font.color.rgb = RGBColor(0x6E, 0x8B, 0xA0) if lit \
                    else RGBColor(0x3A, 0x4A, 0x58)
                run.font.italic = True
            elif not lit:
                run.font.color.rgb = RGBColor(0x45, 0x57, 0x66)
            else:
                run.font.color.rgb = RGBColor(0xE8, 0xEE, 0xF2)
        box.name = "code"
        self._transition(s, "morph" if highlight is not None else "fade")
        self.notes(s, note)
        return s

    def figure_slide(self, title, image, caption="", kicker="", note="",
                     morph=False):
        s = self._new(WHITE)
        self._chrome(s, title, kicker)
        cap_h = 0.72 if caption else 0.0
        area_h = BODY_H - cap_h - 0.10
        if os.path.exists(image):
            from PIL import Image        # pillow ships with anaconda
            try:
                iw, ih = Image.open(image).size
                ar = iw / ih
            except Exception:
                ar = 2.2
            w = min(CW, area_h * ar)
            h = w / ar
            if h > area_h:
                h, w = area_h, area_h * ar
            pic = s.shapes.add_picture(image, Inches(ML + (CW - w) / 2),
                                       Inches(BODY_Y + (area_h - h) / 2),
                                       width=Inches(w), height=Inches(h))
            pic.name = "figure"
        else:
            self.rect(s, ML, BODY_Y, CW, area_h, fill=LIGHTBG, name="figure")
            self.text(s, ML, BODY_Y + area_h / 2 - 0.2, CW, 0.4,
                      f"[figure missing: {os.path.basename(image)}]",
                      size=14, color=ACCENT, align=PP_ALIGN.CENTER,
                      name="figmiss")
        if caption:
            self.text(s, ML, BODY_Y + area_h + 0.06, CW, cap_h, caption,
                      size=13, color=NEUTRAL, align=PP_ALIGN.CENTER,
                      name="figcap")
        self._transition(s, "morph" if morph else "fade")
        self.notes(s, note)
        return s

    def breakfix_slide(self, stage, title, paras=None, bullets=None,
                       code=None, device="", note="", prompt=""):
        """One beat of the guided diagnosis. *stage* labels the band."""
        s = self._new(WHITE)
        band = ACCENT if stage.lower().startswith(("symptom", "break")) \
            else PRIMARY if stage.lower().startswith(("evidence", "diagnos")) \
            else GREENHL
        self.rect(s, 0, 0, W, 0.60, fill=band, name="bfband")
        self.text(s, ML, 0.15, CW * 0.6, 0.32,
                  f"BREAK & FIX  ·  {stage.upper()}",
                  size=14, color=WHITE, bold=True, name="bfstage")
        self.text(s, ML + CW * 0.6, 0.15, CW * 0.4, 0.32,
                  f"Chapter {self.chapter}", size=13,
                  color=RGBColor(0xF2, 0xD7, 0xD5), align=PP_ALIGN.RIGHT,
                  name="bfch")

        y = 0.92
        self.text(s, ML, y, CW, 0.72, title,
                  size=fit_pt(title, CW, 0.70, ideal=27, floor=18),
                  color=NEARBLACK, bold=True, name="bftitle")
        y += 0.86

        # Centre short content in the space left over. Without this a
        # two-line symptom sits under the title with three inches of white
        # below it, which reads as an unfinished slide rather than a spare one.
        region_top, region_bot = y, 6.95
        if prompt:
            ph = 1.10
            y = region_top + max(0.0, (region_bot - region_top - ph) / 2) - 0.6
            self.rect(s, ML, y, CW, 0.92, fill=LIGHTBG, name="promptbg")
            self.text(s, ML + 0.3, y + 0.22, CW - 0.6, 0.52, prompt,
                      size=22, color=PRIMARY, bold=True,
                      align=PP_ALIGN.CENTER, name="prompt")
            y += 1.10
        elif not code and (paras or bullets):
            # Nudged, not centred -- the same cap the code panel gets. Six
            # beats of one diagnosis that each start at a different height
            # make the click between them read as a jump.
            est = _est_h(list(paras or []) + list(bullets or []), CW, 19)
            if est < (region_bot - region_top) * 0.62:
                y = region_top + min(0.55, max(
                    0.0, (region_bot - region_top - est) / 2.6))

        avail = 6.95 - y
        if code:
            lines = code.split("\n")
            size = fit_pt("\n".join(lines), CW - 0.6, avail - 0.55,
                          ideal=15, floor=9)
            needed = 0.36 + 0.30 + len(lines) * 1.24 * size / 72.0
            ph = max(1.10, min(avail, needed))
            # Same cap as terminal_slide: hang the panel from the title
            # rather than optically centring it, so consecutive beats of one
            # diagnosis do not shuffle the evidence up and down the screen.
            if ph < avail - 0.6:
                y += min(0.42, (avail - ph) / 2.6)
            self.rect(s, ML, y, CW, ph, fill=CODEBG, name="bfcodebg")
            self.rect(s, ML, y, CW, 0.36, fill=RGBColor(0x27, 0x35, 0x43),
                      name="bfcodebar")
            self.text(s, ML + 0.22, y + 0.05, CW - 0.44, 0.26, device or "IOS",
                      size=12, color=RGBColor(0x9F, 0xB3, 0xC8), bold=True,
                      name="bfcodettl")
            box = s.shapes.add_textbox(Inches(ML + 0.22), Inches(y + 0.46),
                                       Inches(CW - 0.44), Inches(ph - 0.52))
            tf = box.text_frame
            tf.word_wrap = False
            tf.margin_left = tf.margin_right = Emu(0)
            tf.margin_top = tf.margin_bottom = Emu(0)
            for i, ln in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(0)
                p.line_spacing = 1.0
                r = p.add_run()
                r.text = ln or " "
                r.font.size = Pt(size)
                r.font.name = MONO_FONT
                r.font.color.rgb = RGBColor(0xE8, 0xEE, 0xF2)
            box.name = "bfcode"
        else:
            txt = "\n".join(_flat(p) for p in (paras or [])) + \
                  "\n".join(_flat(b) for b in (bullets or []))
            size = fit_pt(txt, CW - 0.2, avail, ideal=19, floor=13)
            yy = y
            if paras:
                self.text(s, ML, yy, CW, avail, list(paras), size=size,
                          name="bfbody", space_after=10)
                yy += _est_h(paras, CW, size) + 0.12
            if bullets:
                self._bullets(s, ML, yy, CW, 6.95 - yy, bullets, size,
                              name="bfbullets")
        self._footer(s)
        self._transition(s, "morph")
        self.notes(s, note)
        return s

    def question_slide(self, n, total, question, answer=None, note=""):
        """A checkpoint question; *answer* None shows the question alone."""
        s = self._new(PRIMARY if answer is None else WHITE)
        if answer is None:
            self.rect(s, 0, 0, W, H, fill=PRIMARY, name="qbg")
            self.text(s, ML, 1.30, CW, 0.34,
                      f"CHECKPOINT  {n} of {total}", size=14,
                      color=RGBColor(0x9F, 0xC5, 0xE8), bold=True, name="qlab")
            self.text(s, ML, 1.95, CW, 3.3, list([question]),
                      size=fit_pt(_flat(question), CW, 3.0, ideal=30, floor=18),
                      color=WHITE, name="qtext")
            self.text(s, ML, 5.90, CW, 0.34, "Discuss · then advance",
                      size=13, color=RGBColor(0x7F, 0xA8, 0xC9), name="qhint")
        else:
            self._chrome(s, f"Checkpoint {n} of {total}", "ANSWER")
            self.text(s, ML, BODY_Y, CW, 1.5, list([question]),
                      size=fit_pt(_flat(question), CW, 1.3, ideal=20, floor=14),
                      color=PRIMARY, bold=True, name="qtext")
            gap = min(1.7, 0.30 + _est_h([question], CW, 19))
            room = BODY_H - gap - 0.1
            asize = fit_pt(_flat(answer), CW - 0.3, room, ideal=19, floor=13)
            # The rule beside the answer marks the answer, so it should be as
            # tall as the answer -- not run to the footer past empty space.
            abar_h = max(0.5, min(room, _est_h([answer], CW - 0.3, asize,
                                              tight=True) + 0.12))
            self.rect(s, ML, BODY_Y + gap, 0.045, abar_h,
                      fill=GREENHL, name="abar")
            self.text(s, ML + 0.28, BODY_Y + gap, CW - 0.28, room,
                      list([answer]), size=asize,
                      color=NEARBLACK, name="atext")
        # The question is a scene change from whatever preceded it, so it
        # cuts in; the answer morphs out of the question it is answering.
        self._transition(s, "fade" if answer is None else "morph")
        self.notes(s, note)
        return s

    def closing_slide(self, next_title: str = "", takeaways=None) -> None:
        s = self._new(self.ground)
        self.rect(s, 0, 0, W, H, fill=self.ground, name="bg")
        self.text(s, ML, 1.15, CW, 0.9, "Where we got to",
                  size=34, color=WHITE, bold=True, name="cltitle")
        if takeaways:
            # Chapter summaries are written as full sentences, so four of
            # them at 19pt overflow. Size to the actual text instead.
            keep = takeaways[:4]
            body = " ".join(_flat(t) for t in keep)
            size = fit_pt(body, CW - 0.4, 3.30, ideal=19, floor=12)
            box = s.shapes.add_textbox(Inches(ML), Inches(2.20),
                                       Inches(CW), Inches(3.35))
            tf = box.text_frame
            tf.word_wrap = True
            for i, t in enumerate(keep):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(9)
                _hang(p, "•  ", size)
                r = p.add_run()
                r.text = "•  "
                r.font.size = Pt(size)
                r.font.color.rgb = RGBColor(0xC8, 0xD6, 0xE5)
                self._runs_into(p, t, size, WHITE, False, None)
            box.name = "cltakeaways"
        # The last chapter has nothing to point forward to, and a closing
        # slide that just stops is a poor end to a semester. The band says
        # what has actually been finished rather than offering sentiment.
        lab, line = ("NEXT", next_title) if next_title else (
            "END OF THE COURSE",
            "34 chapters  ·  all 29 blueprint topics  ·  32 labs")
        self.rect(s, ML, 5.80, CW, 0.90,
                  fill=RGBColor(*[max(0, c - 22) for c in self.ground]),
                  name="nextbg")
        self.text(s, ML + 0.30, 5.96, CW - 0.6, 0.30, lab,
                  size=12, color=RGBColor(0xB8, 0xC9, 0xDA), bold=True,
                  name="nextlab")
        self.text(s, ML + 0.30, 6.24, CW - 0.6, 0.36, line,
                  size=19, color=WHITE, bold=True, name="nexttitle")
        self._transition(s, "fade")

    # -- output ------------------------------------------------------------

    def save(self, path: str) -> str:
        os.makedirs(os.path.dirname(path), exist_ok=True)
        self.prs.save(path)
        return path

    @property
    def count(self) -> int:
        return len(self.prs.slides._sldIdLst)


# ---------------------------------------------------------------------------

class _Plain:
    """Adapter so a bare string can travel the same path as a Run."""

    def __init__(self, text):
        self.text = text
        self.bold = self.italic = self.mono = self.term = False
        self.chip = None


def _flat(runs) -> str:
    if isinstance(runs, str):
        return runs
    return "".join(getattr(r, "text", "") for r in runs)


def _est_h(paras, width_in: float, size: int, tight: bool = False) -> float:
    """Estimated rendered height of paragraphs, in inches.

    The default is deliberately generous: 0.50 em per character is wider than
    Segoe UI prose actually sets (about 0.44), and the slack absorbs the extra
    line that word wrapping adds. Everything that *chooses a point size* wants
    that slack, because the cost of being wrong is text off the slide.

    ``tight=True`` measures the characters instead, for the few places that
    size a decoration to text already laid out -- the rule beside a checkpoint
    answer, the tint behind a definition. There the cost of slack is a bar
    that runs on past the words it marks, which just looks like a mistake.
    """
    lh = 1.30 * size / 72.0
    lines = 0
    for p in paras:
        t = _flat(p)
        if tight:
            # 4% for the ragged right edge a real wrap leaves.
            need = text_w(t, size) * 1.04
            lines += max(1, -(-int(need * 1000) // int(width_in * 1000)))
        else:
            per_line = max(1, int(width_in / (0.50 * size / 72.0)))
            lines += max(1, -(-len(t) // per_line))
    return lines * lh + 0.10 * len(paras)
