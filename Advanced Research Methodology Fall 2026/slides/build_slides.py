# -*- coding: utf-8 -*-
"""Build the lecture decks from the handout source.

    python build_slides.py            # 33 chapter decks + course overview
    python build_slides.py --part 3   # only Part III
    python build_slides.py --ch 10    # only chapter 10
    python build_slides.py --per-part # one deck per part instead

Content is parsed out of ../parts/*.tex (see texparse.py), so the slides
cannot drift from the printed volume. Colours, wording and chapter numbering
all come from the handout.
"""
import argparse
import os
import re
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

import pptxfx as FX
import texparse as T

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "out")

# ---------------------------------------------- palette (armhandout.sty)
PRIMARY = RGBColor(0x0E, 0x4D, 0x64)
SECONDARY = RGBColor(0x13, 0x71, 0x77)
ACCENT = RGBColor(0xC8, 0x4B, 0x31)
LIGHTBG = RGBColor(0xE8, 0xF6, 0xF3)
DEFBG = RGBColor(0xFD, 0xF2, 0xE9)
NEUTRAL = RGBColor(0x70, 0x7B, 0x7C)
INK = RGBColor(0x1C, 0x28, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAPER = RGBColor(0xFA, 0xFC, 0xFC)
TRAD = {"Empirical": RGBColor(0x1F, 0x61, 0x8D),
        "Design science": RGBColor(0x11, 0x78, 0x64),
        "Qualitative": RGBColor(0x6C, 0x34, 0x83),
        "Formal": RGBColor(0x92, 0x2B, 0x21)}
KIND_COLOR = {"Definition": ACCENT, "Concept": SECONDARY,
              "Watch out": RGBColor(0xB7, 0x95, 0x0B),
              "Example": RGBColor(0x11, 0x7A, 0x65),
              "Regulation": ACCENT, "": SECONDARY}

FONT = "Segoe UI"
FONT_TITLE = "Georgia"          # the handout is set in a serif; titles echo it

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.85)                # side margin
CONTENT_W = W - 2 * M


# --------------------------------------------------------------- helpers
def _txbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    return tb, tf


def _run(p, text, size=18, bold=False, color=INK, font=FONT, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return r


def _rect(slide, x, y, w, h, fill=None, line=None,
          shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    if sh.has_text_frame:
        sh.text_frame.word_wrap = True
    return sh


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, prs, color):
    _rect(slide, 0, 0, prs.slide_width, prs.slide_height, fill=color)


def _footer(slide, prs, left_text, right_text):
    _rect(slide, 0, H - Inches(0.42), W, Pt(2), fill=LIGHTBG)
    tb, tf = _txbox(slide, M, H - Inches(0.38), CONTENT_W, Inches(0.3))
    p = tf.paragraphs[0]
    _run(p, left_text, size=10, color=NEUTRAL)
    tb2, tf2 = _txbox(slide, M, H - Inches(0.38), CONTENT_W, Inches(0.3))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    _run(p2, right_text, size=10, color=NEUTRAL)


def _slide_head(slide, prs, kicker, title, ch):
    """Standard content-slide furniture: rule, kicker, title."""
    _bg(slide, prs, PAPER)
    _rect(slide, 0, 0, W, Inches(0.16), fill=PRIMARY)
    _rect(slide, M, Inches(0.52), Inches(0.6), Pt(3), fill=ACCENT)
    tb, tf = _txbox(slide, M, Inches(0.68), CONTENT_W, Inches(0.32))
    p = tf.paragraphs[0]
    _run(p, kicker.upper(), size=11, bold=True, color=ACCENT)
    # A long section title wraps to two lines; hold it at 30pt and it spills
    # past the block into the body. Step the size down instead of clipping.
    size = 30 if len(title) <= 52 else (26 if len(title) <= 74 else 23)
    tb2, tf2 = _txbox(slide, M, Inches(1.0), CONTENT_W, Inches(0.95))
    p2 = tf2.paragraphs[0]
    _run(p2, title, size=size, bold=True, color=PRIMARY, font=FONT_TITLE)
    return tb2


def _fit(text, budget):
    if len(text) <= budget:
        return text
    return text[:budget].rsplit(" ", 1)[0].rstrip(" ,;:(") + "…"


# The body region every content slide draws into: below the title block,
# above the footer rule. Lists are top-aligned and generously leaded.
# Centring a short list only moves the dead space above it, which reads
# worse than letting the slide breathe underneath; card stacks, which are
# heavier, are centred instead.
BODY_T = Inches(2.0)
BODY_H = Inches(4.5)


def _body(slide, pad_left=Emu(0), pad_right=Emu(0)):
    tb, tf = _txbox(slide, M + pad_left, BODY_T,
                    CONTENT_W - pad_left - pad_right, BODY_H,
                    anchor=MSO_ANCHOR.TOP)
    return tb, tf


def _est_lines(text, chars_per_line):
    """Rough wrapped-line count, for sizing a card to its content."""
    n = 0
    for para in str(text).split("\n"):
        n += max(1, -(-len(para) // max(1, chars_per_line)))
    return n


# ------------------------------------------------------------ slide types
def slide_title(prs, ch, course):
    s = _blank(prs)
    _bg(s, prs, PRIMARY)
    _rect(s, 0, 0, Inches(0.22), H, fill=ACCENT)
    _rect(s, W - Inches(4.6), 0, Inches(4.6), H, fill=SECONDARY)

    tb, tf = _txbox(s, Inches(1.1), Inches(1.5), Inches(7.4), Inches(0.4))
    p = tf.paragraphs[0]
    _run(p, "CHAPTER %d" % ch["number"], size=14, bold=True, color=LIGHTBG)

    tb2, tf2 = _txbox(s, Inches(1.1), Inches(2.0), Inches(7.3), Inches(2.4))
    p2 = tf2.paragraphs[0]
    _run(p2, ch["title"], size=40, bold=True, color=WHITE, font=FONT_TITLE)

    _rect(s, Inches(1.1), Inches(4.6), Inches(1.3), Pt(3), fill=ACCENT)

    tb3, tf3 = _txbox(s, Inches(1.1), Inches(4.9), Inches(7.3), Inches(1.0))
    p3 = tf3.paragraphs[0]
    _run(p3, "Part %d — %s" % (ch["part"], ch["part_title"]),
         size=16, color=LIGHTBG)
    p4 = tf3.add_paragraph()
    _run(p4, course, size=13, color=LIGHTBG, italic=True)

    tb4, tf4 = _txbox(s, W - Inches(4.1), Inches(2.4), Inches(3.1),
                      Inches(3.0))
    p5 = tf4.paragraphs[0]
    _run(p5, "MS IT / PhD IT", size=15, bold=True, color=WHITE)
    p6 = tf4.add_paragraph()
    p6.space_before = Pt(6)
    _run(p6, "Department of Information Technology\nThe Islamia University of "
             "Bahawalpur", size=12, color=LIGHTBG)
    p7 = tf4.add_paragraph()
    p7.space_before = Pt(14)
    _run(p7, "Fall 2026", size=12, color=LIGHTBG)

    FX.transition(s, "fade")
    FX.animate(s, [FX.click(FX.fx(tb, "fade", 400)),
                   FX.click(FX.fx(tb2, "flyBottom", 600)),
                   FX.click(FX.fx(tb3, "fade", 400)),
                   FX.click(FX.fx(tb4, "wipeRight", 500))])
    return s


def slide_outcomes(prs, ch, foot):
    s = _blank(prs)
    _slide_head(s, prs, "Learning outcomes", "By the end of this chapter", ch)
    tb, tf = _body(s)
    n = 0
    for i, o in enumerate(ch["outcomes"][:7]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        _run(p, "▸  ", size=17, bold=True, color=ACCENT)
        # the leading verb is bolded in the handout; keep that emphasis
        m = re.match(r"^(\w+)\s+(.*)$", o)
        if m:
            _run(p, m.group(1) + " ", size=17, bold=True, color=PRIMARY)
            _run(p, _fit(m.group(2), 150), size=17, color=INK)
        else:
            _run(p, _fit(o, 160), size=17, color=INK)
        n += 1
    _footer(s, prs, foot, "Outcomes")
    FX.transition(s, "push")
    FX.animate(s, FX.cascade(tb, n, "wipeUp", 350))
    return s


def slide_keyterms(prs, ch, foot):
    s = _blank(prs)
    _slide_head(s, prs, "Key terms", "The vocabulary of this chapter", ch)
    terms = ch["keyterms"][:18]
    if not terms:
        return None
    cols, x0 = 3, M
    cw = (CONTENT_W - Inches(0.4)) / cols
    per = (len(terms) + cols - 1) // cols
    boxes = []
    for c in range(cols):
        chunk = terms[c * per:(c + 1) * per]
        if not chunk:
            continue
        tb, tf = _txbox(s, x0 + c * (cw + Inches(0.2)), BODY_T, cw, BODY_H)
        for i, term in enumerate(chunk):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(11)
            _run(p, "▪  ", size=14, color=ACCENT)
            _run(p, _fit(term, 40), size=14, color=INK)
        boxes.append(tb)
    _footer(s, prs, foot, "Key terms")
    FX.transition(s, "fade")
    FX.animate(s, [FX.click(*[FX.fx(b, "zoom", 400) for b in boxes])])
    return s


def slide_section(prs, ch, title, idx, total, foot):
    s = _blank(prs)
    _bg(s, prs, SECONDARY)
    _rect(s, 0, 0, W, Inches(0.16), fill=ACCENT)
    tb, tf = _txbox(s, M, Inches(2.6), CONTENT_W, Inches(0.4))
    p = tf.paragraphs[0]
    _run(p, "SECTION %d OF %d" % (idx, total), size=13, bold=True,
         color=LIGHTBG)
    tb2, tf2 = _txbox(s, M, Inches(3.05), CONTENT_W - Inches(1.0), Inches(1.8))
    p2 = tf2.paragraphs[0]
    _run(p2, title, size=34, bold=True, color=WHITE, font=FONT_TITLE)
    bar = _rect(s, M, Inches(4.75), Inches(1.6), Pt(4), fill=ACCENT)
    FX.transition(s, "wipe")
    FX.animate(s, [FX.click(FX.fx(tb, "fade", 300),
                            FX.fx(tb2, "flyLeft", 500),
                            FX.fx(bar, "wipeRight", 400))])
    return s


def slide_points(prs, ch, sec_title, points, foot, part_label):
    s = _blank(prs)
    _slide_head(s, prs, part_label, _fit(sec_title, 72), ch)
    pts = points[:3]
    # centre the stack so one or two points do not sit under a blank half-slide
    step = Inches(1.55)
    block = step * len(pts) - (step - Inches(1.32))
    y = BODY_T + (BODY_H - block) / 2
    shapes = []
    for pt in pts:
        head = pt["head"] or pt["kind"]
        col = KIND_COLOR.get(pt["kind"], SECONDARY)
        band = _rect(s, M, y, Inches(0.06), Inches(1.32), fill=col)
        tb, tf = _txbox(s, M + Inches(0.28), y, CONTENT_W - Inches(0.28),
                        Inches(1.32))
        if head:
            p = tf.paragraphs[0]
            p.space_after = Pt(4)
            if pt["kind"]:
                _run(p, pt["kind"].upper() + "   ", size=10, bold=True,
                     color=col)
            _run(p, _fit(head, 70), size=17, bold=True, color=PRIMARY)
            p2 = tf.add_paragraph()
        else:
            p2 = tf.paragraphs[0]
        _run(p2, _fit(pt["text"], 300), size=14, color=INK)
        shapes.append((band, tb))
        y += step
    _footer(s, prs, foot, sec_title[:46])
    FX.transition(s, "fade")
    FX.animate(s, [FX.click(FX.fx(b, "wipeUp", 300), FX.fx(t, "fade", 400))
                   for b, t in shapes])
    return s


def slide_traditions(prs, ch, foot):
    tr = ch["traditions"]
    if not tr:
        return None
    s = _blank(prs)
    _slide_head(s, prs, "Four traditions",
                "How this lands in %s" % _fit(tr["topic"], 46), ch)
    quads = list(tr["quadrants"].items())
    bw = (CONTENT_W - Inches(0.35)) / 2
    bh = Inches(1.95)
    boxes = []
    for i, (lab, txt) in enumerate(quads[:4]):
        cx = M + (i % 2) * (bw + Inches(0.35))
        cy = Inches(2.05) + (i // 2) * (bh + Inches(0.3))
        col = TRAD.get(lab, SECONDARY)
        card = _rect(s, cx, cy, bw, bh, fill=WHITE, line=LIGHTBG)
        _rect(s, cx, cy, bw, Pt(4), fill=col)
        tb, tf = _txbox(s, cx + Inches(0.22), cy + Inches(0.22),
                        bw - Inches(0.44), bh - Inches(0.4))
        p = tf.paragraphs[0]
        p.space_after = Pt(5)
        _run(p, lab.upper(), size=11, bold=True, color=col)
        p2 = tf.add_paragraph()
        _run(p2, _fit(txt, 250), size=12.5, color=INK)
        boxes.append((card, tb))
    _footer(s, prs, foot, "Four traditions")
    FX.transition(s, "fade")
    FX.animate(s, [FX.click(FX.fx(c, "zoom", 350), FX.fx(t, "fade", 400))
                   for c, t in boxes])
    return s


def slide_worked(prs, ch, foot):
    w = ch["worked"]
    if not w:
        return None
    s = _blank(prs)
    _slide_head(s, prs, "Worked example", _fit(w["title"], 68), ch)

    body_text = _fit(w["text"], 300) if w["text"] else ""
    items = [_fit(i, 150) for i in w["items"][:5]]
    # Size the card to its content. Some worked boxes are built round a table
    # we cannot lift, leaving only a sentence or two -- a fixed-height card
    # then renders as a large empty panel.
    lines = (_est_lines(body_text, 105) if body_text else 0)
    lines += sum(_est_lines(i, 108) for i in items)
    card_h = Inches(0.55) + Inches(0.30) * lines + Inches(0.10) * len(items)
    card_h = max(Inches(1.35), min(card_h, Inches(4.0)))
    card_y = BODY_T + (BODY_H - Inches(0.5) - card_h) / 2

    card = _rect(s, M, card_y, CONTENT_W, card_h, fill=DEFBG)
    _rect(s, M, card_y, Inches(0.07), card_h, fill=ACCENT)
    tb, tf = _txbox(s, M + Inches(0.35), card_y + Inches(0.22),
                    CONTENT_W - Inches(0.7), card_h - Inches(0.44))
    n = 0
    if body_text:
        p = tf.paragraphs[0]
        p.space_after = Pt(10)
        _run(p, body_text, size=15, color=INK)
        n += 1
    for it in items:
        p = tf.add_paragraph() if n else tf.paragraphs[0]
        p.space_after = Pt(7)
        _run(p, "•  ", size=14, bold=True, color=ACCENT)
        _run(p, it, size=14, color=INK)
        n += 1
    note, ntf = _txbox(s, M, card_y + card_h + Inches(0.16),
                       CONTENT_W, Inches(0.4))
    _run(ntf.paragraphs[0],
         "Constructed for teaching — the numbers are invented, as in the "
         "handout.", size=11, color=NEUTRAL, italic=True)
    _footer(s, prs, foot, "Worked example")
    FX.transition(s, "fade")
    steps = [FX.click(FX.fx(card, "wipeUp", 400))]
    steps += FX.cascade(tb, n, "fade", 350)
    steps.append(FX.click(FX.fx(note, "fade", 300)))
    FX.animate(s, steps)
    return s


def slide_pitfall(prs, ch, foot):
    items = ch.get("pitfall_items") or (
        [ch["pitfall"]] if ch["pitfall"] else [])
    if not items:
        return None
    s = _blank(prs)
    _slide_head(s, prs, "Pitfalls",
                _fit(ch.get("pitfall_title") or "Common mistakes", 68), ch)
    tb, tf = _body(s, pad_left=Inches(0.1), pad_right=Inches(0.1))
    n = 0
    for i, it in enumerate(items[:6]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(15)
        _run(p, "✖  ", size=15, bold=True, color=ACCENT)
        m = re.match(r"^(.{3,60}?[.:])\s+(.*)$", it)
        if m:
            _run(p, m.group(1) + " ", size=15, bold=True, color=PRIMARY)
            _run(p, _fit(m.group(2), 140), size=15, color=INK)
        else:
            _run(p, _fit(it, 170), size=15, color=INK)
        n += 1
    _footer(s, prs, foot, "Pitfalls")
    FX.transition(s, "fade")
    FX.animate(s, FX.cascade(tb, n, "flyLeft", 350))
    return s


def slide_reg(prs, ch, foot):
    regs = ch["regboxes"][:3]
    if not regs:
        return None
    s = _blank(prs)
    _slide_head(s, prs, "What the regulations say",
                "Rules that bind you, quoted", ch)
    step = Inches(1.55)
    block = step * len(regs) - (step - Inches(1.35))
    y = BODY_T + (BODY_H - block) / 2
    shapes = []
    for r in regs:
        card = _rect(s, M, y, CONTENT_W, Inches(1.35), fill=WHITE, line=ACCENT)
        _rect(s, M, y, Inches(0.07), Inches(1.35), fill=ACCENT)
        tb, tf = _txbox(s, M + Inches(0.3), y + Inches(0.16),
                        CONTENT_W - Inches(0.6), Inches(1.05))
        p = tf.paragraphs[0]
        p.space_after = Pt(4)
        _run(p, _fit(r["clause"], 70), size=12, bold=True, color=ACCENT)
        p2 = tf.add_paragraph()
        _run(p2, _fit(r["text"], 250), size=13, color=INK)
        shapes.append((card, tb))
        y += step
    _footer(s, prs, foot, "Regulations")
    FX.transition(s, "fade")
    FX.animate(s, [FX.click(FX.fx(c, "wipeRight", 350), FX.fx(t, "fade", 400))
                   for c, t in shapes])
    return s


def slide_summary(prs, ch, foot):
    if not ch["summary"]:
        return None
    s = _blank(prs)
    _bg(s, prs, LIGHTBG)
    _rect(s, 0, 0, W, Inches(0.16), fill=PRIMARY)
    _rect(s, M, Inches(0.52), Inches(0.6), Pt(3), fill=ACCENT)
    tb0, tf0 = _txbox(s, M, Inches(0.68), CONTENT_W, Inches(0.32))
    _run(tf0.paragraphs[0], "SUMMARY", size=11, bold=True, color=ACCENT)
    tb1, tf1 = _txbox(s, M, Inches(1.0), CONTENT_W, Inches(0.8))
    _run(tf1.paragraphs[0], "What to take away", size=30, bold=True,
         color=PRIMARY, font=FONT_TITLE)
    tb, tf = _body(s)
    n = 0
    for i, it in enumerate(ch["summary"][:7]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        _run(p, "%d. " % (i + 1), size=15, bold=True, color=ACCENT)
        _run(p, _fit(it, 165), size=15, color=INK)
        n += 1
    _footer(s, prs, foot, "Summary")
    FX.transition(s, "push")
    FX.animate(s, FX.cascade(tb, n, "wipeUp", 320))
    return s


def slide_questions(prs, ch, foot):
    qs = ch["reviewq"] or ch["checkpoint"]
    if not qs:
        return None
    s = _blank(prs)
    _slide_head(s, prs, "Review", "Questions to think with", ch)
    tb, tf = _body(s)
    n = 0
    for i, q in enumerate(qs[:5]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(18)
        _run(p, "Q%d  " % (i + 1), size=15, bold=True, color=ACCENT)
        _run(p, _fit(q, 190), size=15, color=INK)
        n += 1
    _footer(s, prs, foot, "Review questions")
    FX.transition(s, "fade")
    FX.animate(s, FX.cascade(tb, n, "fade", 350))
    return s


def slide_close(prs, ch, nxt, foot):
    s = _blank(prs)
    _bg(s, prs, PRIMARY)
    _rect(s, 0, 0, Inches(0.22), H, fill=ACCENT)
    tb, tf = _txbox(s, Inches(1.4), Inches(2.7), Inches(9.0), Inches(1.0))
    _run(tf.paragraphs[0], "End of Chapter %d" % ch["number"], size=34,
         bold=True, color=WHITE, font=FONT_TITLE)
    tb2, tf2 = _txbox(s, Inches(1.4), Inches(3.9), Inches(9.5), Inches(1.4))
    p = tf2.paragraphs[0]
    if nxt:
        _run(p, "Next: ", size=17, color=LIGHTBG)
        _run(p, "Ch. %d — %s" % (nxt["number"], nxt["title"]),
             size=17, bold=True, color=WHITE)
    p2 = tf2.add_paragraph()
    p2.space_before = Pt(12)
    _run(p2, "Read the chapter in the handout before the next session.",
         size=13, color=LIGHTBG, italic=True)
    FX.transition(s, "fade")
    FX.animate(s, [FX.click(FX.fx(tb, "flyBottom", 500)),
                   FX.click(FX.fx(tb2, "fade", 400))])
    return s


# ------------------------------------------------------------- assembly
def build_chapter_slides(prs, ch, nxt, course):
    foot = "Ch. %d · %s" % (ch["number"], _fit(ch["title"], 54))
    part_label = "Part %d · Ch. %d" % (ch["part"], ch["number"])
    slide_title(prs, ch, course)
    slide_outcomes(prs, ch, foot)
    slide_keyterms(prs, ch, foot)
    secs = [x for x in ch["sections"] if x["points"]]
    for i, sec in enumerate(secs, 1):
        slide_section(prs, ch, sec["title"], i, len(secs), foot)
        pts = sec["points"]
        for k in range(0, len(pts), 3):
            slide_points(prs, ch, sec["title"], pts[k:k + 3], foot, part_label)
    slide_worked(prs, ch, foot)
    slide_pitfall(prs, ch, foot)
    slide_reg(prs, ch, foot)
    slide_traditions(prs, ch, foot)
    slide_summary(prs, ch, foot)
    slide_questions(prs, ch, foot)
    slide_close(prs, ch, nxt, foot)


def new_prs():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def overview_deck(chapters, course):
    prs = new_prs()
    s = _blank(prs)
    _bg(s, prs, PRIMARY)
    _rect(s, 0, 0, Inches(0.22), H, fill=ACCENT)
    tb, tf = _txbox(s, Inches(1.2), Inches(2.2), Inches(10.0), Inches(2.0))
    _run(tf.paragraphs[0], "Advanced Research Methodology", size=44,
         bold=True, color=WHITE, font=FONT_TITLE)
    tb2, tf2 = _txbox(s, Inches(1.2), Inches(4.1), Inches(10.0), Inches(1.6))
    _run(tf2.paragraphs[0], "MS IT / PhD IT  ·  Fall 2026", size=18,
         color=LIGHTBG)
    p = tf2.add_paragraph()
    p.space_before = Pt(10)
    _run(p, "Department of Information Technology · The Islamia "
            "University of Bahawalpur", size=13, color=LIGHTBG)
    p2 = tf2.add_paragraph()
    p2.space_before = Pt(14)
    _run(p2, "Required by PhD Degree Regulations-2024 §10(iv)(a)",
         size=12, color=LIGHTBG, italic=True)
    FX.transition(s, "fade")
    FX.animate(s, [FX.click(FX.fx(tb, "flyBottom", 600)),
                   FX.click(FX.fx(tb2, "fade", 450))])

    by_part = {}
    for ch in chapters:
        by_part.setdefault(ch["part"], []).append(ch)
    for pno in sorted(by_part):
        chs = by_part[pno]
        s = _blank(prs)
        _slide_head(s, prs, "Part %d" % pno,
                    T.PART_TITLES.get(pno, ""), chs[0])
        tb, tf = _txbox(s, M, Inches(2.05), CONTENT_W, Inches(4.4))
        n = 0
        for i, ch in enumerate(chs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(8)
            _run(p, "%2d. " % ch["number"], size=15, bold=True, color=ACCENT)
            _run(p, _fit(ch["title"], 90), size=15, color=INK)
            n += 1
        _footer(s, prs, "Course overview", "Part %d" % pno)
        FX.transition(s, "push")
        FX.animate(s, FX.cascade(tb, n, "wipeUp", 260))
    return prs


def safe(name):
    return re.sub(r"[^A-Za-z0-9]+", "_", name).strip("_")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--ch", type=int, help="build only this chapter number")
    ap.add_argument("--part", type=int, help="build only this part")
    ap.add_argument("--per-part", action="store_true",
                    help="one deck per part instead of per chapter")
    ap.add_argument("--no-overview", action="store_true")
    args = ap.parse_args()

    if not os.path.isdir(OUT):
        os.makedirs(OUT)
    course = "Advanced Research Methodology · Fall 2026"
    nums = T.chapter_numbers()
    if not nums:
        print("WARNING: no .aux found; chapter numbers will be 0. "
              "Build the handout first.", file=sys.stderr)

    chapters = []
    for fn in T.chapter_files():
        ch = T.parse_chapter(os.path.join(T.PARTS, fn), nums)
        if ch:
            chapters.append(ch)
    chapters.sort(key=lambda c: c["number"])

    sel = chapters
    if args.ch:
        sel = [c for c in chapters if c["number"] == args.ch]
    if args.part:
        sel = [c for c in sel if c["part"] == args.part]

    written = []
    if args.per_part:
        by_part = {}
        for ch in sel:
            by_part.setdefault(ch["part"], []).append(ch)
        for pno in sorted(by_part):
            prs = new_prs()
            chs = by_part[pno]
            for i, ch in enumerate(chs):
                nxt = chs[i + 1] if i + 1 < len(chs) else None
                build_chapter_slides(prs, ch, nxt, course)
            path = os.path.join(OUT, "ARM_Part%d_%s.pptx"
                                % (pno, safe(T.PART_TITLES.get(pno, ""))[:40]))
            prs.save(path)
            written.append((path, len(prs.slides._sldIdLst)))
    else:
        for ch in sel:
            prs = new_prs()
            i = chapters.index(ch)
            nxt = chapters[i + 1] if i + 1 < len(chapters) else None
            build_chapter_slides(prs, ch, nxt, course)
            path = os.path.join(OUT, "ARM_Ch%02d_%s.pptx"
                                % (ch["number"], safe(ch["title"])[:44]))
            prs.save(path)
            written.append((path, len(prs.slides._sldIdLst)))

    if not args.no_overview and not args.ch:
        prs = overview_deck(chapters, course)
        path = os.path.join(OUT, "ARM_00_Course_Overview.pptx")
        prs.save(path)
        written.append((path, len(prs.slides._sldIdLst)))

    total = sum(n for _, n in written)
    for path, n in written:
        print("  %-64s %3d slides" % (os.path.basename(path), n))
    print("\n%d deck(s), %d slides -> %s" % (len(written), total, OUT))


if __name__ == "__main__":
    main()
