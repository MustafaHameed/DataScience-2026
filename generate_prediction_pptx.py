"""
Generate Prediction/ML PowerPoint Presentation
Topic: Predictive Modeling with Machine Learning
Author: Dr. Mustafa Hameed
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def add_slide(prs, layout_index, title_text, content_items=None, subtitle_text=None):
    """Add a slide with title and optional bullet content."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title_text
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.size = Pt(28)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 51, 102)
    
    # Set subtitle if layout supports it
    if subtitle_text and layout_index == 0:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.text = subtitle_text
                for para in ph.text_frame.paragraphs:
                    para.font.size = Pt(18)
                    para.font.color.rgb = RGBColor(80, 80, 80)
    
    # Add content items as bullets
    if content_items and layout_index == 1:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                tf = ph.text_frame
                tf.clear()
                for i, item in enumerate(content_items):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = item
                    p.font.size = Pt(18)
                    p.space_after = Pt(6)
                    if item.startswith("  "):
                        p.level = 1
                        p.font.size = Pt(16)
                break
    
    return slide


def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.alignment = alignment
    if color:
        p.font.color.rgb = RGBColor(*color)
    return txBox


def add_code_box(slide, left, top, width, height, code_text):
    """Add a styled code box."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
    shape.line.color.rgb = RGBColor(180, 180, 180)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    
    for i, line in enumerate(code_text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(30, 30, 30)


def add_table_slide(prs, title_text, headers, rows, col_widths=None):
    """Add a slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Title
    add_text_box(slide, 0.5, 0.2, 9, 0.6, title_text, font_size=28, bold=True, color=(0, 51, 102))
    
    num_rows = len(rows) + 1
    num_cols = len(headers)
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9.0)
    height = Inches(0.4 * num_rows)
    
    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table
    
    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    
    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
    
    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(230, 240, 250)
    
    return slide


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    add_slide(prs, 0,
              "Predictive Modeling with Machine Learning",
              subtitle_text="Predicting Student Performance from LMS Data\n\nDr. Mustafa Hameed")
    
    # ==========================================
    # SLIDE 2: Learning Objectives
    # ==========================================
    add_slide(prs, 1,
              "Learning Objectives",
              [
                  "Understand the concept of predictive modeling",
                  "Learn the machine learning workflow (7 steps)",
                  "Perform Exploratory Data Analysis (EDA)",
                  "Build Linear Regression and Random Forest models",
                  "Evaluate models using RMSE, MAE, and R²",
                  "Interpret results and variable importance",
              ])
    
    # ==========================================
    # SLIDE 3: What is Predictive Modeling?
    # ==========================================
    add_slide(prs, 1,
              "What is Predictive Modeling?",
              [
                  "Uses historical data to forecast future/unknown outcomes",
                  "In education: predict student performance from behavior data",
                  "We solve a REGRESSION problem (predicting a number)",
                  "",
                  "Applications in Education:",
                  "  • Early warning systems for at-risk students",
                  "  • Personalized learning recommendations",
                  "  • Resource allocation and support planning",
              ])
    
    # ==========================================
    # SLIDE 4: Supervised vs Unsupervised
    # ==========================================
    add_table_slide(prs,
                    "Supervised vs Unsupervised Learning",
                    ["Aspect", "Supervised", "Unsupervised"],
                    [
                        ["Labels", "Known outcomes (grades)", "No known outcomes"],
                        ["Goal", "Predict an output", "Find patterns/clusters"],
                        ["Examples", "Regression, Classification", "Clustering, PCA"],
                        ["Our Case", "✓ Supervised Regression", "—"],
                    ],
                    col_widths=[2.0, 3.5, 3.5])
    
    # ==========================================
    # SLIDE 5: ML Workflow
    # ==========================================
    add_slide(prs, 1,
              "The Machine Learning Workflow",
              [
                  "Step 1: Load Data",
                  "Step 2: Explore (EDA) — distributions, correlations",
                  "Step 3: Prepare Data — standardize, handle issues",
                  "Step 4: Split — training (80%) & testing (20%)",
                  "Step 5: Build Model(s) — train on training data",
                  "Step 6: Evaluate — RMSE, MAE, R² on test data",
                  "Step 7: Interpret & Compare Results",
              ])
    
    # ==========================================
    # SLIDE 6: The Dataset
    # ==========================================
    add_slide(prs, 1,
              "The Dataset: Student LMS Engagement",
              [
                  "285 students, 12 numeric variables",
                  "Source: github.com/lamethods/data2 (synthetic data)",
                  "",
                  "Features (11 predictors):",
                  "  • Freq_Course_View, Freq_Lecture_View",
                  "  • Freq_Forum_Consume, Freq_Forum_Contribute",
                  "  • Regularity (4 variables for consistency)",
                  "  • Session_Count, Total_Duration, Active_Days",
                  "",
                  "Target: Final_Grade (continuous numeric)",
              ])
    
    # ==========================================
    # SLIDE 7: Variables Table
    # ==========================================
    add_table_slide(prs,
                    "Dataset Variables",
                    ["Variable", "Description", "Type"],
                    [
                        ["Freq_Course_View", "Frequency of course page views", "Predictor"],
                        ["Freq_Lecture_View", "Frequency of lecture views", "Predictor"],
                        ["Freq_Forum_Consume", "Frequency of reading forums", "Predictor"],
                        ["Freq_Forum_Contribute", "Frequency of forum posts", "Predictor"],
                        ["Session_Count", "Number of login sessions", "Predictor"],
                        ["Total_Duration", "Total time on LMS", "Predictor"],
                        ["Active_Days", "Number of active days", "Predictor"],
                        ["Final_Grade", "Student final grade", "Target"],
                    ],
                    col_widths=[3.0, 3.5, 1.5])
    
    # ==========================================
    # SLIDE 8: Loading Data in R
    # ==========================================
    slide = add_slide(prs, 1,
                      "Step 1: Load the Data in R",
                      [
                          "Use the rio package for easy data import",
                          "glimpse() shows structure and types",
                      ])
    add_code_box(slide, 0.8, 3.5, 8.4, 1.8,
                 'library(tidyverse)\nlibrary(rio)\n\nlms <- import("https://github.com/lamethods/data2/raw/main/lms/lms.csv")\nglimpse(lms)')
    
    # ==========================================
    # SLIDE 9: EDA - Summary Stats
    # ==========================================
    slide = add_slide(prs, 1,
                      "Step 2: Exploratory Data Analysis",
                      [
                          "Before modeling, understand the data:",
                          "  • Check distributions and ranges",
                          "  • Look for missing values",
                          "  • Examine correlations with target",
                      ])
    add_code_box(slide, 0.8, 3.8, 8.4, 1.5,
                 'library(skimr)\nskim(lms)                         # Summary statistics\ncor(lms)[, "Final_Grade"]         # Correlations with target')
    
    # ==========================================
    # SLIDE 10: Correlation
    # ==========================================
    add_table_slide(prs,
                    "Correlation Strength Guide",
                    ["Correlation (|r|)", "Strength"],
                    [
                        ["0.00 – 0.19", "Very Weak"],
                        ["0.20 – 0.39", "Weak"],
                        ["0.40 – 0.59", "Moderate"],
                        ["0.60 – 0.79", "Strong"],
                        ["0.80 – 1.00", "Very Strong"],
                    ],
                    col_widths=[4.5, 4.5])
    
    # ==========================================
    # SLIDE 11: Correlation Heatmap Code
    # ==========================================
    slide = add_slide(prs, 1,
                      "EDA: Correlation Heatmap",
                      [
                          "Visualize all pairwise correlations",
                          "Red = positive, Blue = negative correlation",
                      ])
    add_code_box(slide, 0.8, 3.2, 8.4, 2.5,
                 'cor_matrix <- cor(lms)\ncor_df <- as.data.frame(as.table(cor_matrix))\nnames(cor_df) <- c("Var1", "Var2", "Correlation")\n\nggplot(cor_df, aes(x=Var1, y=Var2, fill=Correlation)) +\n  geom_tile(color="white") +\n  scale_fill_gradient2(low="blue", high="red", mid="white") +\n  geom_text(aes(label=round(Correlation,2)), size=2.5)')
    
    # ==========================================
    # SLIDE 12: Scatter Plots
    # ==========================================
    slide = add_slide(prs, 1,
                      "EDA: Feature vs Target Scatter Plots",
                      [
                          "Visualize relationship between each feature and Final_Grade",
                          "Look for linear vs non-linear patterns",
                      ])
    add_code_box(slide, 0.8, 3.2, 8.4, 2.3,
                 'lms_long <- lms %>%\n  pivot_longer(cols = -Final_Grade,\n               names_to = "Feature", values_to = "Value")\n\nggplot(lms_long, aes(x=Value, y=Final_Grade)) +\n  geom_point(alpha=0.3, color="steelblue") +\n  geom_smooth(method="lm", se=FALSE, color="red") +\n  facet_wrap(~Feature, scales="free_x", ncol=4)')
    
    # ==========================================
    # SLIDE 13: Standardization
    # ==========================================
    add_slide(prs, 1,
              "Step 3: Data Preparation — Standardization",
              [
                  "Transform features to mean=0, std=1:",
                  "",
                  "  z = (x − μ) / σ",
                  "",
                  "Why standardize?",
                  "  • Features are on different scales",
                  "  • Some algorithms (KNN, SVM) are scale-sensitive",
                  "  • Ensures fair comparison of feature importance",
              ])
    
    # ==========================================
    # SLIDE 14: Standardization Code
    # ==========================================
    slide = add_slide(prs, 1,
                      "Standardization in R",
                      [
                          "Use scale() to standardize predictors",
                          "Keep the target variable (Final_Grade) unchanged",
                      ])
    add_code_box(slide, 0.8, 3.2, 8.4, 2.0,
                 '# Separate predictors and target\npredictors <- lms %>% select(-Final_Grade)\ntarget <- lms$Final_Grade\n\n# Standardize and recombine\nlms_scaled <- as.data.frame(scale(predictors))\nlms_scaled$Final_Grade <- target')
    
    # ==========================================
    # SLIDE 15: VIF
    # ==========================================
    add_slide(prs, 1,
              "Multicollinearity Check (VIF)",
              [
                  "VIF = Variance Inflation Factor",
                  "Detects when predictors are too correlated with each other",
                  "",
                  "  VIF = 1 → No collinearity",
                  "  VIF 1–5 → Moderate (acceptable)",
                  "  VIF > 5 → High (consider removing feature)",
                  "  VIF > 10 → Very high (definitely problematic)",
              ])
    
    # ==========================================
    # SLIDE 16: Train/Test Split
    # ==========================================
    add_slide(prs, 1,
              "Step 4: Train/Test Split",
              [
                  "Divide data into two separate sets:",
                  "",
                  "  Training Set (80%) — model learns patterns here",
                  "  Test Set (20%) — evaluate on unseen data",
                  "",
                  "Why split?",
                  "  • Prevents overfitting (memorizing training data)",
                  "  • Gives honest estimate of real-world performance",
                  "  • set.seed(123) ensures reproducibility",
              ])
    
    # ==========================================
    # SLIDE 17: Split Code
    # ==========================================
    slide = add_slide(prs, 1,
                      "Train/Test Split in R",
                      [
                          "Use rsample package for easy splitting",
                      ])
    add_code_box(slide, 0.8, 2.8, 8.4, 2.2,
                 'library(rsample)\nset.seed(123)  # For reproducibility\n\nsplit <- initial_split(lms_scaled, prop = 0.80)\ntrain_data <- training(split)   # ~228 rows\ntest_data  <- testing(split)    # ~57 rows\n\ncat("Training:", nrow(train_data), "rows")\ncat("Testing: ", nrow(test_data), "rows")')
    
    # ==========================================
    # SLIDE 18: Linear Regression Theory
    # ==========================================
    add_slide(prs, 1,
              "Model 1: Linear Regression",
              [
                  "Fits a straight-line relationship:",
                  "",
                  "  ŷ = β₀ + β₁x₁ + β₂x₂ + ... + βₚxₚ",
                  "",
                  "Strengths:",
                  "  • Easy to interpret (coefficients = effects)",
                  "  • Fast to train",
                  "",
                  "Weaknesses:",
                  "  • Assumes linear relationships only",
                  "  • Sensitive to outliers",
              ])
    
    # ==========================================
    # SLIDE 19: LR Code
    # ==========================================
    slide = add_slide(prs, 1,
                      "Linear Regression: R Code",
                      [
                          "Train the model, then predict on test set",
                      ])
    add_code_box(slide, 0.8, 2.8, 8.4, 2.5,
                 '# Train Linear Regression\nmodel_lm <- lm(Final_Grade ~ ., data = train_data)\nsummary(model_lm)\n\n# Predict on test set\npred_lm <- predict(model_lm, newdata = test_data)\n\n# View results\nhead(data.frame(Actual = test_data$Final_Grade,\n                Predicted = round(pred_lm, 3)))')
    
    # ==========================================
    # SLIDE 20: Random Forest Theory
    # ==========================================
    add_slide(prs, 1,
              "Model 2: Random Forest",
              [
                  "An ENSEMBLE method that builds many decision trees:",
                  "",
                  "  1. Create 500 trees on random data subsets (bagging)",
                  "  2. Each split uses random feature subsets",
                  "  3. Final prediction = average of all trees",
                  "",
                  "Strengths: Non-linear, robust, variable importance",
                  "Weaknesses: Less interpretable, slower to train",
              ])
    
    # ==========================================
    # SLIDE 21: RF Code
    # ==========================================
    slide = add_slide(prs, 1,
                      "Random Forest: R Code",
                      [
                          "Train with 500 trees, then predict",
                      ])
    add_code_box(slide, 0.8, 2.8, 8.4, 2.8,
                 'library(randomForest)\nset.seed(123)\n\n# Train Random Forest\nmodel_rf <- randomForest(\n  Final_Grade ~ .,\n  data = train_data,\n  ntree = 500,\n  importance = TRUE\n)\n\n# Predict on test set\npred_rf <- predict(model_rf, newdata = test_data)')
    
    # ==========================================
    # SLIDE 22: Evaluation Metrics
    # ==========================================
    add_slide(prs, 1,
              "Step 6: Evaluation Metrics",
              [
                  "RMSE — Root Mean Squared Error",
                  "  • Average prediction error (lower = better)",
                  "  • Penalizes large errors more",
                  "",
                  "MAE — Mean Absolute Error",
                  "  • Average absolute error (lower = better)",
                  "  • Less sensitive to outliers",
                  "",
                  "R² — Coefficient of Determination",
                  "  • % of variance explained (higher = better)",
                  "  • R²=0.5 means model explains 50% of variance",
              ])
    
    # ==========================================
    # SLIDE 23: Metric Formulas
    # ==========================================
    add_table_slide(prs,
                    "Evaluation Metric Formulas",
                    ["Metric", "Formula", "Goal"],
                    [
                        ["RMSE", "√(Σ(yᵢ - ŷᵢ)² / n)", "Lower = Better"],
                        ["MAE", "Σ|yᵢ - ŷᵢ| / n", "Lower = Better"],
                        ["R²", "1 - SS_res / SS_tot", "Higher = Better"],
                    ],
                    col_widths=[2.0, 4.5, 2.5])
    
    # ==========================================
    # SLIDE 24: Evaluation Code
    # ==========================================
    slide = add_slide(prs, 1,
                      "Evaluation in R",
                      [
                          "Use yardstick package for metrics",
                      ])
    add_code_box(slide, 0.8, 2.5, 8.4, 3.5,
                 'library(yardstick)\n\n# Linear Regression metrics\ncat("=== Linear Regression ===")\ncat("RMSE:", rmse_vec(test_data$Final_Grade, pred_lm))\ncat("MAE: ", mae_vec(test_data$Final_Grade, pred_lm))\ncat("R²:  ", rsq_vec(test_data$Final_Grade, pred_lm))\n\n# Random Forest metrics\ncat("=== Random Forest ===")\ncat("RMSE:", rmse_vec(test_data$Final_Grade, pred_rf))\ncat("MAE: ", mae_vec(test_data$Final_Grade, pred_rf))\ncat("R²:  ", rsq_vec(test_data$Final_Grade, pred_rf))')
    
    # ==========================================
    # SLIDE 25: Predicted vs Actual
    # ==========================================
    slide = add_slide(prs, 1,
                      "Visualization: Predicted vs Actual",
                      [
                          "Points close to the diagonal = good predictions",
                          "Points far from diagonal = prediction errors",
                      ])
    add_code_box(slide, 0.8, 3.2, 8.4, 2.5,
                 'viz_data <- data.frame(\n  Actual = rep(test_data$Final_Grade, 2),\n  Predicted = c(pred_lm, pred_rf),\n  Model = rep(c("LR", "RF"), each = nrow(test_data))\n)\nggplot(viz_data, aes(x=Actual, y=Predicted, color=Model)) +\n  geom_point(alpha=0.6) +\n  geom_abline(slope=1, intercept=0, linetype="dashed") +\n  facet_wrap(~Model)')
    
    # ==========================================
    # SLIDE 26: Residual Plot
    # ==========================================
    add_slide(prs, 1,
              "Visualization: Residual Plots",
              [
                  "Residual = Actual − Predicted",
                  "",
                  "Good model: residuals randomly scattered around 0",
                  "Bad model: patterns visible in residuals",
                  "",
                  "What to check:",
                  "  • No fan shape (heteroscedasticity)",
                  "  • No curve (non-linearity missed)",
                  "  • Centered around zero line",
              ])
    
    # ==========================================
    # SLIDE 27: Variable Importance
    # ==========================================
    slide = add_slide(prs, 1,
                      "Visualization: Variable Importance",
                      [
                          "Random Forest tells us which features matter most",
                          "%IncMSE: how much error increases without this feature",
                          "Higher = more important",
                      ])
    add_code_box(slide, 0.8, 3.5, 8.4, 1.5,
                 '# Variable importance plot\nvarImpPlot(model_rf,\n           main = "Variable Importance")\n\n# Or extract as a data frame\nimportance(model_rf)')
    
    # ==========================================
    # SLIDE 28: Model Comparison
    # ==========================================
    add_table_slide(prs,
                    "When to Use Each Model",
                    ["Scenario", "Recommended Model"],
                    [
                        ["Need interpretable results", "Linear Regression"],
                        ["Need best accuracy", "Random Forest"],
                        ["Small dataset", "Linear Regression"],
                        ["Non-linear relationships", "Random Forest"],
                        ["Feature importance needed", "Random Forest"],
                        ["Speed is important", "Linear Regression"],
                    ],
                    col_widths=[5.0, 4.0])
    
    # ==========================================
    # SLIDE 29: Key Findings
    # ==========================================
    add_slide(prs, 1,
              "Key Findings & Interpretation",
              [
                  "Both models predict Final Grade from LMS engagement",
                  "",
                  "Most Important Features (typically):",
                  "  • Session_Count — number of login sessions",
                  "  • Active_Days — days with LMS activity",
                  "  • Total_Duration — time spent on platform",
                  "  • Freq_Lecture_View — lecture viewing frequency",
                  "",
                  "Practical Implication:",
                  "  Students who engage more frequently and consistently",
                  "  tend to achieve higher final grades",
              ])
    
    # ==========================================
    # SLIDE 30: Key Definitions
    # ==========================================
    add_table_slide(prs,
                    "Key Definitions",
                    ["Term", "Definition"],
                    [
                        ["Supervised Learning", "Learning from labeled examples (known outcomes)"],
                        ["Regression", "Predicting a continuous numeric value"],
                        ["Training Set", "Data the model learns from (80%)"],
                        ["Test Set", "Held-out data for evaluation (20%)"],
                        ["Overfitting", "Model memorizes noise instead of patterns"],
                        ["RMSE", "Root Mean Squared Error — average error"],
                        ["MAE", "Mean Absolute Error — avg absolute error"],
                        ["R²", "Proportion of variance explained (0 to 1)"],
                    ],
                    col_widths=[2.5, 6.5])
    
    # ==========================================
    # SLIDE 31: Practice Exercises
    # ==========================================
    add_slide(prs, 1,
              "Practice Exercises",
              [
                  "1. Change split to 70/30 — how do results change?",
                  "2. Remove 2 least important features — does R² improve?",
                  "3. Calculate RMSE by hand for first 5 test observations",
                  "4. Interpret Linear Regression coefficients — which",
                  "   features have positive vs negative effects?",
                  "5. Create a bar chart comparing metrics for both models",
                  "6. Research: What is k-fold cross-validation?",
              ])
    
    # ==========================================
    # SLIDE 32: Summary
    # ==========================================
    add_slide(prs, 1,
              "Summary",
              [
                  "1. Loaded 285-student LMS dataset (12 variables)",
                  "2. Explored data with EDA (distributions, correlations)",
                  "3. Standardized features for fair comparison",
                  "4. Split 80/20 into training and testing",
                  "5. Built Linear Regression and Random Forest models",
                  "6. Evaluated with RMSE, MAE, and R²",
                  "7. Identified key engagement features for success",
                  "",
                  "Next Steps: Cross-validation, more algorithms (SVM, KNN)",
              ])
    
    # ==========================================
    # SLIDE 33: References
    # ==========================================
    add_slide(prs, 1,
              "References & Resources",
              [
                  "Dataset: github.com/lamethods/data2 (Synthetic LMS Data)",
                  "",
                  "Textbook Chapter:",
                  "  Estrellado et al. — Learning Analytics Methods",
                  "  Chapter 3: Prediction",
                  "  lamethods.org/book2/chapters/ch03-prediction/",
                  "",
                  "R Packages Used:",
                  "  tidyverse, rio, skimr, randomForest,",
                  "  rsample, yardstick, performance",
              ])
    
    # Save
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                               "R in DataScience", "Prediction_ML_Presentation.pptx")
    # Also try current dir if R in DataScience doesn't exist at that level
    output_dir = r"e:\Documents\Github\DataScience-2026\R in DataScience"
    output_path = os.path.join(output_dir, "Prediction_ML_Presentation.pptx")
    
    prs.save(output_path)
    print(f"Presentation saved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
